"""
app.py — FulcrumQ Deck Converter
Streamlit UI for convert_deck.py

Run:
    streamlit run app.py
"""

import base64
import contextlib
import io
import json
import re
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path

import streamlit as st
import streamlit.components.v1 as components

# ── Page config ───────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="FulcrumQ Deck Converter",
    page_icon=str(Path(__file__).parent / "logo package" / "PNG" / "monogram_default.png"),
    layout="wide",
)

BASE_DIR = Path(__file__).parent
sys.path.insert(0, str(BASE_DIR))
import convert_deck as cd

import os as _os
# Load .env from project dir if present (so ANTHROPIC_API_KEY doesn't need to be set manually)
_env_file = BASE_DIR / ".env"
if _env_file.exists():
    for _line in _env_file.read_text().splitlines():
        _line = _line.strip()
        if _line and not _line.startswith("#") and "=" in _line:
            _k, _v = _line.split("=", 1)
            _os.environ.setdefault(_k.strip(), _v.strip())
try:
    from slide_vision import classify_deck_vision, CONFIDENCE_THRESHOLD as _VIS_THRESHOLD
    _HAS_VISION = True
except ImportError:
    _HAS_VISION = False

from lxml import etree as _ET

# ── Semantic Ingest tab ───────────────────────────────────────────────────────
try:
    from tabs.semantic_ingest_tab import render_semantic_ingest_tab
    _HAS_SEMANTIC_INGEST = True
except ImportError:
    _HAS_SEMANTIC_INGEST = False

# ── 4-layer pipeline ──────────────────────────────────────────────────────────
try:
    from classifier    import classify_deck    as _pl_classify
    from schema_engine import run_schema_engine as _pl_schema
    from renderer      import render_deck       as _pl_render
    from qa_validator  import validate_render   as _pl_qa
    _PIPELINE_AVAILABLE = True
    _PIPELINE_MASTER = BASE_DIR / "FulcrumQ_Blank_Template.pptx"
    _PIPELINE_RULES  = BASE_DIR / "layout_rules.json"
except ImportError as _pie:
    _PIPELINE_AVAILABLE = False
    _PIPELINE_MASTER    = None
    _PIPELINE_RULES     = None

# ── Audit helpers ─────────────────────────────────────────────────────────────
_BRAND_COLORS = {
    "765FFF","917FFF","AD9FFF","C8BFFF","E9E4FF",
    "1D1D1D","3B3B3B","585858","767676","7A828D",
    "B2BBCA","D0D7DF","F2F2F2","FFFFFF","00C27A","FF2E88",
    "FFB547","60BDBC","EA3323","281A42",
}
_BRAND_FONTS = {"Segoe UI", "Arial"}
_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _audit_source(pptx_bytes: bytes) -> dict:
    """Scan source PPTX and return colors + fonts not covered by any mapping."""
    unmapped_colors: dict[str, int] = {}
    unmapped_fonts:  dict[str, int] = {}
    with zipfile.ZipFile(io.BytesIO(pptx_bytes)) as z:
        slide_keys = [k for k in z.namelist()
                      if re.match(r"ppt/slides/slide\d+\.xml$", k)]
        for key in slide_keys:
            root = _ET.fromstring(z.read(key))
            for srgb in root.iter(f"{{{_NS_A}}}srgbClr"):
                v = srgb.get("val", "").upper()
                if (v and v not in cd.COLOR_REMAP
                        and not cd._is_red_hue(v)
                        and v not in _BRAND_COLORS):
                    unmapped_colors[v] = unmapped_colors.get(v, 0) + 1
            for latin in root.iter(f"{{{_NS_A}}}latin"):
                face = latin.get("typeface", "")
                if (face and face not in ("+mj-lt", "+mn-lt", "")
                        and face not in cd.FONT_MAP
                        and face not in _BRAND_FONTS):
                    unmapped_fonts[face] = unmapped_fonts.get(face, 0) + 1
    return {"colors": unmapped_colors, "fonts": unmapped_fonts}


def _brand_check(pptx_bytes: bytes) -> list[dict]:
    """Scan a PPTX per-slide for off-brand colors, fonts, and CEO Works text.
    Returns list of {slide, colors: {hex: count}, fonts: {face: count}, names: int}.
    """
    issues = []
    with zipfile.ZipFile(io.BytesIO(pptx_bytes)) as z:
        slide_keys = sorted(
            [k for k in z.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", k)],
            key=lambda k: int(re.search(r"\d+", k.split("/")[-1]).group()),
        )
        for i, key in enumerate(slide_keys, 1):
            root = _ET.fromstring(z.read(key))
            bad_colors: dict[str, int] = {}
            bad_fonts:  dict[str, int] = {}
            name_count = 0

            for srgb in root.iter(f"{{{_NS_A}}}srgbClr"):
                v = srgb.get("val", "").upper()
                if v and v not in _BRAND_COLORS:
                    bad_colors[v] = bad_colors.get(v, 0) + 1

            for latin in root.iter(f"{{{_NS_A}}}latin"):
                face = latin.get("typeface", "")
                if face and face not in ("+mj-lt", "+mn-lt", "") and face not in _BRAND_FONTS:
                    bad_fonts[face] = bad_fonts.get(face, 0) + 1

            for t_el in root.iter(f"{{{_NS_A}}}t"):
                if t_el.text and cd._CEO_WORKS_RE.search(t_el.text):
                    name_count += 1

            if bad_colors or bad_fonts or name_count:
                issues.append({
                    "slide":  i,
                    "colors": bad_colors,
                    "fonts":  bad_fonts,
                    "names":  name_count,
                })
    return issues


def _render_slide_thumbs(pptx_path: Path, width_px: int = 320) -> list[bytes]:
    """Render PPTX slides → PNG thumbnails via LibreOffice + PyMuPDF.
    Returns [] silently if LibreOffice times out or is unavailable."""
    try:
        import fitz
    except ImportError:
        return []
    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        try:
            subprocess.run(
                ["soffice", "--headless", "--convert-to", "pdf",
                 "--outdir", str(tmp_path), str(pptx_path)],
                capture_output=True, timeout=240,
            )
        except (subprocess.TimeoutExpired, FileNotFoundError):
            return []
        pdf_path = tmp_path / (pptx_path.stem + ".pdf")
        if not pdf_path.exists():
            return []
        doc = fitz.open(str(pdf_path))
        thumbs = []
        for page in doc:
            scale = width_px / page.rect.width
            mat   = fitz.Matrix(scale, scale)
            pix   = page.get_pixmap(matrix=mat, alpha=False)
            thumbs.append(pix.tobytes("png"))
        doc.close()
        return thumbs


# ── Brand CSS ─────────────────────────────────────────────────────────────────
st.markdown("""
<style>
/* Segoe UI (Windows system font) → Inter as the closest web-safe substitute
   on macOS. Weights loaded: 400 Regular, 600 SemiBold, 700 Bold, 900 Black.
   Arial is the brand body font — always available as a system font. */
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;600;700;900&display=swap');

/* ── Brand tokens — sourced directly from FulcrumQ color guide ── */
:root {
    /* Typography — mirrors master PPTX font assignment:
       Headlines : Segoe UI  (Inter on macOS, exact metric match)
       Body      : Arial     (system font, no substitute needed)
       Mono      : Courier New / Fira Code for code/log blocks        */
    --fq-font-head: 'Segoe UI', 'Inter', Arial, sans-serif;
    --fq-font-body: Arial, sans-serif;
    --fq-font-mono: 'Courier New', 'Fira Code', monospace;
    /* Pivot Purple + tints (15-20% usage) */
    --fq-purple:        #765FFF;
    --fq-tint1:         #917FFF;
    --fq-tint2:         #AD9FFF;
    --fq-tint3:         #C8BFFF;

    /* Guiding Grey + shades (75-80% base, dark UI) */
    --fq-grey:          #1D1D1D;
    --fq-grey1:         #3B3B3B;
    --fq-grey2:         #585858;
    --fq-grey3:         #767676;

    /* Named greys */
    --fq-grey-mid:      #7A828D;
    --fq-light-grey:    #B2BBCA;
    --fq-soft-grey:     #D0D7DF;

    /* Secondary accents (3-5% combined MAX) */
    --fq-vector-dp:     #281A42;
    --fq-green:         #00C27A;
    --fq-magenta:       #FF2E88;
    --fq-amber:         #FFB547;
    --fq-teal:          #60BDBC;

    /* White */
    --fq-white:         #FFFFFF;

    /* Utility */
    --fq-panel:     rgba(255,255,255,0.03);
    --fq-border:    rgba(118,95,255,0.22);
}

/* ── Base ── */
html, body,
[data-testid="stAppViewContainer"],
[data-testid="stMain"] > div {
    background: var(--fq-grey) !important;
    color: #e8e8f0;
    font-family: var(--fq-font-body);
    font-size: 12px;
    font-weight: 400;
}
[data-testid="stHeader"],
[data-testid="stToolbar"] { display: none !important; }
#MainMenu, footer { visibility: hidden !important; }

.block-container {
    padding-top: 1.5rem !important;
    max-width: 1080px;
}

/* ── Keyframes ── */
@keyframes gradient-pan {
    0%   { background-position: 0% 50%; }
    50%  { background-position: 100% 50%; }
    100% { background-position: 0% 50%; }
}
@keyframes fold-in {
    0% {
        opacity: 0;
        transform: perspective(700px) rotateX(-18deg) translateY(-16px) scale(0.97);
    }
    100% {
        opacity: 1;
        transform: perspective(700px) rotateX(0deg) translateY(0) scale(1);
    }
}
@keyframes fold-in-left {
    0% {
        opacity: 0;
        transform: perspective(700px) rotateY(20deg) translateX(-20px) scale(0.97);
    }
    100% {
        opacity: 1;
        transform: perspective(700px) rotateY(0deg) translateX(0) scale(1);
    }
}
@keyframes border-pulse {
    0%, 100% {
        border-color: rgba(118,95,255,0.35);
        box-shadow: 0 0 0 0 rgba(118,95,255,0);
    }
    50% {
        border-color: rgba(118,95,255,0.8);
        box-shadow: 0 0 22px 4px rgba(118,95,255,0.18);
    }
}
@keyframes fade-up {
    from { opacity: 0; transform: translateY(14px); }
    to   { opacity: 1; transform: translateY(0); }
}
@keyframes ripple-out {
    0%   { transform: scale(0); opacity: 0.5; }
    100% { transform: scale(4.5); opacity: 0; }
}
@keyframes shimmer {
    0%   { background-position: -400px 0; }
    100% { background-position: 400px 0; }
}
@keyframes glow-pulse {
    /* Pivot Purple #765FFF → Signal Magenta #FF2E88 */
    0%, 100% { filter: drop-shadow(0 0 8px rgba(118,95,255,0.6)); }
    50%       { filter: drop-shadow(0 0 18px rgba(255,46,136,0.7)); }
}

/* ── Hero header ── */
.fq-header {
    background: linear-gradient(135deg,
        var(--fq-grey)   0%,
        var(--fq-grey1)  45%,
        var(--fq-purple) 100%
    );
    background-size: 300% 300%;
    animation: gradient-pan 8s ease infinite;
    border-radius: 16px;
    padding: 2rem 2.5rem 1.75rem;
    margin-bottom: 0.25rem;
    border: 1px solid var(--fq-border);
    position: relative;
    overflow: hidden;
}
.fq-header::before {
    content: '';
    position: absolute;
    top: -30%; right: -5%;
    width: 320px; height: 320px;
    /* Signal Magenta #FF2E88 soft glow */
    background: radial-gradient(circle, rgba(255,46,136,0.10) 0%, transparent 68%);
    pointer-events: none;
}
.fq-header::after {
    content: '';
    position: absolute;
    bottom: -40%; left: 20%;
    width: 240px; height: 240px;
    background: radial-gradient(circle, rgba(0,194,122,0.08) 0%, transparent 68%);
    pointer-events: none;
}
.fq-header h1,
.fq-header h1 * {
    font-family: 'Segoe UI', 'Inter', Arial, sans-serif !important;
    font-size: 2rem;
    font-weight: 700;
    margin: 0 0 0.35rem;
    letter-spacing: -0.025em;
    background: linear-gradient(90deg, var(--fq-white) 25%, var(--fq-tint2) 75%, var(--fq-tint1) 100%);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    background-clip: text;
}
/* Hide Streamlit auto-injected heading anchor */
.fq-header h1 a, .fq-header h1 svg { display: none !important; }

.fq-header p {
    font-family: var(--fq-font-body);
    font-size: 0.82rem;
    font-weight: 400;
    color: rgba(255,255,255,0.48);
    margin: 0;
    letter-spacing: 0.01em;
}
.fq-steps {
    display: flex;
    gap: 0.5rem;
    margin-top: 1rem;
    flex-wrap: wrap;
}
.fq-step-pill {
    font-family: var(--fq-font-head);
    font-weight: 700;
    background: rgba(255,255,255,0.07);
    border: 1px solid rgba(255,255,255,0.12);
    border-radius: 20px;
    padding: 3px 10px;
    font-size: 10px;
    color: rgba(255,255,255,0.52);
    letter-spacing: 0.06em;
    transition: background 0.2s, color 0.2s, border-color 0.2s;
    cursor: default;
}
.fq-step-pill:hover {
    background: rgba(118,95,255,0.18);
    border-color: rgba(118,95,255,0.45);
    color: var(--fq-tint3);
}

/* ── Tabs ── */
[data-testid="stTabs"] [data-baseweb="tab-list"] {
    background: transparent !important;
    border-bottom: 1px solid var(--fq-border) !important;
    gap: 4px;
}
[data-testid="stTabs"] [data-baseweb="tab"],
[data-testid="stTabs"] [data-baseweb="tab"] *,
[data-testid="stTabs"] [data-baseweb="tab-list"] button,
[data-testid="stTabs"] [data-baseweb="tab-list"] button * {
    background: transparent !important;
    color: var(--fq-grey3) !important;
    font-family: Arial, sans-serif !important;
    font-size: 0.78rem !important;
    font-weight: 700 !important;
    letter-spacing: 0.08em !important;
    text-transform: uppercase !important;
    padding: 0.55rem 1.3rem !important;
    border-radius: 8px 8px 0 0 !important;
    border: none !important;
    transition: color 0.2s, background 0.2s !important;
    position: relative;
}
[data-testid="stTabs"] [data-baseweb="tab"]::after {
    content: '';
    position: absolute;
    bottom: 0; left: 50%; right: 50%;
    height: 2px;
    background: var(--fq-purple);
    transition: left 0.25s, right 0.25s;
    border-radius: 2px 2px 0 0;
}
[data-testid="stTabs"] [data-baseweb="tab"]:hover::after {
    left: 20%; right: 20%;
}
[data-testid="stTabs"] [aria-selected="true"]::after {
    left: 0; right: 0;
}
[data-testid="stTabs"] [data-baseweb="tab"]:hover {
    color: rgba(255,255,255,0.85) !important;
    background: rgba(118,95,255,0.08) !important;
}
[data-testid="stTabs"] [aria-selected="true"] {
    color: var(--fq-tint2) !important;
    background: rgba(118,95,255,0.1) !important;
}

/* ── Upload zone ── */
[data-testid="stFileUploader"] {
    animation: fold-in 0.45s cubic-bezier(0.23,1,0.32,1) both;
}
[data-testid="stFileUploader"] section {
    background: var(--fq-panel) !important;
    border: 2px dashed rgba(118,95,255,0.38) !important;
    border-radius: 14px !important;
    padding: 2.2rem 1.5rem !important;
    animation: border-pulse 4s ease-in-out infinite !important;
    transition: background 0.25s, transform 0.25s, border-color 0.25s !important;
}
[data-testid="stFileUploader"] section:hover {
    background: rgba(118,95,255,0.05) !important;
    transform: scale(1.012) !important;
    border-color: rgba(118,95,255,0.7) !important;
    animation: none !important;
    box-shadow: 0 0 28px rgba(118,95,255,0.15) !important;
}
[data-testid="stFileUploaderDropzone"] {
    color: var(--fq-tint2) !important;
}
[data-testid="stFileUploaderDropzoneInstructions"] span {
    color: rgba(255,255,255,0.5) !important;
    font-size: 0.82rem !important;
}

/* ── Primary button — Guiding Grey → Pivot Purple ── */
[data-testid="stBaseButton-primary"],
.stButton > button[kind="primary"] {
    background: linear-gradient(135deg, var(--fq-grey) 0%, var(--fq-purple) 100%) !important;
    border: none !important;
    border-radius: 9px !important;
    color: var(--fq-white) !important;
    font-family: Arial, sans-serif !important;
    font-weight: 700 !important;
    font-size: 0.82rem !important;
    letter-spacing: 0.05em !important;
    text-transform: uppercase !important;
    padding: 0.6rem 1.6rem !important;
    position: relative;
    overflow: hidden;
    transition: transform 0.18s, box-shadow 0.18s !important;
    will-change: transform;
}
[data-testid="stBaseButton-primary"]:hover,
.stButton > button[kind="primary"]:hover {
    transform: translateY(-2px) !important;
    box-shadow: 0 10px 30px rgba(118,95,255,0.45) !important;
}
[data-testid="stBaseButton-primary"]:active,
.stButton > button[kind="primary"]:active {
    transform: translateY(0) scale(0.96) !important;
    box-shadow: 0 4px 12px rgba(118,95,255,0.3) !important;
}

/* ── Download button — Anchor Green #00C27A ── */
.stDownloadButton > button {
    background: transparent !important;
    border: 1.5px solid var(--fq-green) !important;
    color: var(--fq-green) !important;
    border-radius: 9px !important;
    font-family: Arial, sans-serif !important;
    font-weight: 700 !important;
    font-size: 0.82rem !important;
    letter-spacing: 0.05em !important;
    text-transform: uppercase !important;
    transition: background 0.22s, transform 0.18s, box-shadow 0.22s !important;
}
.stDownloadButton > button * {
    font-family: Arial, sans-serif !important;
    color: inherit !important;
}
.stDownloadButton > button:hover {
    background: rgba(0,194,122,0.08) !important;
    transform: translateY(-2px) !important;
    box-shadow: 0 8px 24px rgba(0,194,122,0.25) !important;
}
.stDownloadButton > button:active {
    transform: translateY(0) scale(0.96) !important;
}

/* ── Alerts ── */
[data-testid="stAlert"][data-baseweb="notification"] {
    background: rgba(0,194,122,0.06) !important;
    border: 1px solid rgba(0,194,122,0.26) !important;
    border-radius: 10px !important;
    animation: fold-in 0.38s ease both !important;
}
[data-testid="stAlert"][kind="error"] {
    background: rgba(255,46,136,0.06) !important;
    border: 1px solid rgba(255,46,136,0.26) !important;
}

/* ── Log terminal ── */
.log-box {
    background: #0a0a12;
    border: 1px solid rgba(118,95,255,0.18);
    border-radius: 12px;
    padding: 1.1rem 1.3rem;
    font-family: var(--fq-font-mono);
    font-size: 11px;
    white-space: pre-wrap;
    max-height: 360px;
    overflow-y: auto;
    color: #b8b8cc;
    line-height: 1.65;
    animation: fold-in 0.45s cubic-bezier(0.23,1,0.32,1) both;
    position: relative;
}
.log-box::before {
    content: '● ● ●';
    display: block;
    color: rgba(118,95,255,0.4);
    font-size: 10px;
    letter-spacing: 3px;
    margin-bottom: 0.6rem;
    font-family: monospace;
}
.log-box::-webkit-scrollbar { width: 5px; }
.log-box::-webkit-scrollbar-track { background: transparent; }
.log-box::-webkit-scrollbar-thumb {
    background: rgba(118,95,255,0.3);
    border-radius: 3px;
}

/* ── Section label ── */
.section-label {
    font-family: var(--fq-font-head);
    font-size: 9px;
    font-weight: 700;
    text-transform: uppercase;
    letter-spacing: 0.14em;
    color: var(--fq-grey3);
    margin: 1.5rem 0 0.5rem;
}

/* ── Color swatches ── */
.swatch-row {
    display: flex;
    align-items: center;
    gap: 8px;
    margin: 3px 0;
    font-family: var(--fq-font-mono);
    font-size: 11px;
    color: rgba(255,255,255,0.62);
    padding: 5px 8px;
    border-radius: 7px;
    transition: background 0.22s, transform 0.22s, color 0.22s;
    cursor: default;
}
.swatch-row:hover {
    background: rgba(118,95,255,0.1);
    transform: translateX(6px);
    color: rgba(255,255,255,0.88);
}
.swatch {
    width: 19px; height: 19px;
    border-radius: 50%;
    display: inline-block;
    flex-shrink: 0;
    border: 1px solid rgba(255,255,255,0.1);
    transition: transform 0.22s, box-shadow 0.22s;
}
.swatch-row:hover .swatch {
    transform: scale(1.35);
    box-shadow: 0 2px 10px rgba(0,0,0,0.5);
}
.arrow {
    color: rgba(118,95,255,0.55);
    font-weight: 700;
    transition: color 0.22s, transform 0.22s;
    display: inline-block;
}
.swatch-row:hover .arrow {
    color: var(--fq-tint1);
    transform: translateX(2px);
}

/* ── Master badge ── */
.master-badge {
    display: inline-flex;
    align-items: center;
    gap: 6px;
    padding: 5px 14px;
    border-radius: 20px;
    font-family: var(--fq-font-head);
    font-size: 11px;
    font-weight: 700;
    letter-spacing: 0.03em;
    animation: fold-in-left 0.4s ease both;
    transition: box-shadow 0.2s;
    cursor: default;
}
.master-ok {
    background: rgba(0,194,122,0.08);
    border: 1px solid rgba(0,194,122,0.28);
    color: var(--fq-green);
}
.master-ok:hover { box-shadow: 0 0 16px rgba(0,194,122,0.18); }
.master-err {
    background: rgba(255,46,136,0.07);
    border: 1px solid rgba(255,46,136,0.28);
    color: var(--fq-magenta);
}

/* ── Horizontal rule ── */
hr { border-color: rgba(118,95,255,0.12) !important; }

/* ── Spinner ── */
[data-testid="stSpinner"] svg { color: var(--fq-purple) !important; }

/* ── Scheme fill table ── */
.scheme-row {
    display: flex;
    align-items: center;
    gap: 10px;
    margin: 4px 0;
    font-family: var(--fq-font-mono);
    font-size: 11px;
    color: rgba(255,255,255,0.62);
    padding: 5px 8px;
    border-radius: 7px;
    transition: background 0.22s, transform 0.22s;
    cursor: default;
}
.scheme-row:hover {
    background: rgba(118,95,255,0.1);
    transform: translateX(6px);
    color: rgba(255,255,255,0.9);
}

/* ── Hero download box ── */
.fq-hero-dl {
    display: inline-flex;
    flex-direction: column;
    align-items: flex-end;
    gap: 4px;
    padding: 14px 18px;
    border: 1.5px solid rgba(0,194,122,0.55);
    border-radius: 12px;
    background: rgba(0,194,122,0.07);
    text-decoration: none;
    transition: background 0.2s, border-color 0.2s, box-shadow 0.2s;
    cursor: pointer;
    white-space: nowrap;
    box-shadow: 0 0 0 3px rgba(0,194,122,0.12), 0 0 28px rgba(0,194,122,0.18);
}
.fq-hero-dl:hover {
    background: rgba(0,194,122,0.13);
    border-color: rgba(0,194,122,0.85);
    box-shadow: 0 0 0 3px rgba(0,194,122,0.2), 0 0 40px rgba(0,194,122,0.3);
}
.fq-hero-dl-label {
    font-family: Arial, sans-serif;
    font-size: 9px;
    font-weight: 700;
    letter-spacing: 0.12em;
    text-transform: uppercase;
    color: #00C27A;
}
.fq-hero-dl-name {
    font-family: 'Segoe UI', 'Inter', Arial, sans-serif;
    font-size: 1rem;
    font-weight: 700;
    color: #ffffff;
    line-height: 1.2;
}
.fq-hero-dl-ext {
    font-family: Arial, sans-serif;
    font-size: 9px;
    font-weight: 400;
    color: rgba(255,255,255,0.45);
    letter-spacing: 0.06em;
}

/* ── Pivot line — white → Pivot Purple ── */
.fq-pivot-line {
    height: 3px;
    background: linear-gradient(90deg,
        rgba(255,255,255,0.9) 0%,
        rgba(255,255,255,0.6) 30%,
        var(--fq-purple)     66%,
        var(--fq-tint1)      100%
    );
    border-radius: 2px;
    margin: 1.25rem 0;
}

/* ── Icon grid cells ── */
.icon-cell {
    background: rgba(118,95,255,0.04);
    border: 1px solid rgba(118,95,255,0.14);
    border-radius: 10px;
    padding: 8px;
    display: flex;
    align-items: center;
    justify-content: center;
    transition: background 0.18s, border-color 0.18s, box-shadow 0.18s;
    cursor: default;
}
.icon-cell:hover {
    background: rgba(118,95,255,0.10);
    border-color: rgba(118,95,255,0.40);
    box-shadow: 0 0 14px rgba(118,95,255,0.18);
}

/* ── Unmapped audit table ── */
.audit-table { width: 100%; border-collapse: collapse; font-size: 11px; font-family: monospace; }
.audit-table th {
    font-family: Arial, sans-serif; font-size: 9px; font-weight: 700;
    letter-spacing: 0.12em; text-transform: uppercase;
    color: var(--fq-grey3); padding: 4px 8px; text-align: left;
    border-bottom: 1px solid var(--fq-border);
}
.audit-table td { padding: 4px 8px; color: rgba(255,255,255,0.65); }
.audit-table tr:hover td { background: rgba(118,95,255,0.06); }
.audit-dot { display:inline-block; width:12px; height:12px; border-radius:50%;
             vertical-align:middle; margin-right:6px; border:1px solid rgba(255,255,255,0.15); }

/* ── Slide preview thumbnails ── */
.slide-thumb {
    border-radius: 8px;
    border: 1px solid var(--fq-border);
    overflow: hidden;
    transition: box-shadow 0.18s, transform 0.18s;
    cursor: default;
}
.slide-thumb:hover {
    transform: translateY(-3px);
    box-shadow: 0 10px 28px rgba(0,0,0,0.5);
}
.slide-thumb img { display: block; width: 100%; }
.slide-thumb-num {
    font-family: Arial, sans-serif; font-size: 9px;
    color: rgba(255,255,255,0.3); text-align: center; padding: 3px 0 4px;
    background: rgba(0,0,0,0.25);
}

/* ── Typography specimen ── */
.type-row {
    display: flex; align-items: baseline; gap: 16px;
    padding: 14px 0; border-bottom: 1px solid rgba(118,95,255,0.10);
}
.type-meta {
    width: 200px; flex-shrink: 0;
    font-family: Arial, sans-serif; font-size: 9px;
    color: var(--fq-grey3); line-height: 1.8;
}
.type-meta strong { display: block; color: rgba(255,255,255,0.6); font-size: 10px; margin-bottom: 2px; }
.type-sample { flex: 1; line-height: 1.2; }

/* ── Logo tile ── */
.logo-tile {
    border-radius: 14px;
    padding: 28px 20px 16px;
    display: flex;
    flex-direction: column;
    align-items: center;
    gap: 12px;
    border: 1px solid rgba(255,255,255,0.07);
    transition: transform 0.2s cubic-bezier(0.23,1,0.32,1),
                box-shadow 0.2s cubic-bezier(0.23,1,0.32,1);
}
.logo-tile:hover {
    transform: translateY(-3px);
    box-shadow: 0 12px 32px rgba(0,0,0,0.4);
}
.logo-variant-label {
    font-family: var(--fq-font-head);
    font-size: 9px;
    font-weight: 700;
    text-transform: uppercase;
    letter-spacing: 0.12em;
    margin-top: 4px;
}

/* ── Expanders — force dark background so content is readable ── */
[data-testid="stExpander"] {
    background: #111118 !important;
    border: 1px solid var(--fq-border) !important;
    border-radius: 10px !important;
}
[data-testid="stExpander"] summary {
    color: rgba(255,255,255,0.75) !important;
    font-family: Arial, sans-serif !important;
    font-size: 0.82rem !important;
}
[data-testid="stExpander"] summary:hover {
    color: #ffffff !important;
}
[data-testid="stExpander"] > div[data-testid="stExpanderDetails"] {
    background: #111118 !important;
}

/* ── Kill Streamlit red — radio dot only ── */
[data-testid="stTabs"] [data-baseweb="tab-highlight"] {
    background-color: #FF2E88 !important;
}
/* Target only the inner filled dot of the selected radio button */
[data-baseweb="radio"] [data-checked="true"] > div > div,
[data-baseweb="radio"] [role="radio"][aria-checked="true"] > div > div {
    background-color: #FF2E88 !important;
    border-color: #FF2E88 !important;
}
</style>
""", unsafe_allow_html=True)

# ── JS: ripple, tilt, fold ────────────────────────────────────────────────────
components.html("""
<script>
(function() {
    const doc = window.parent.document;

    function attachRipple(btn) {
        if (btn._fqRipple) return;
        btn._fqRipple = true;
        btn.style.position = 'relative';
        btn.style.overflow = 'hidden';
        btn.addEventListener('pointerdown', function(e) {
            const rect = btn.getBoundingClientRect();
            const d = Math.max(btn.offsetWidth, btn.offsetHeight);
            const r = doc.createElement('span');
            r.style.cssText = [
                'position:absolute','border-radius:50%',
                `width:${d}px`,`height:${d}px`,
                'background:rgba(255,255,255,0.22)',
                'transform:scale(0)','animation:ripple-out 0.58s linear',
                'pointer-events:none',
                `left:${e.clientX - rect.left - d/2}px`,
                `top:${e.clientY - rect.top - d/2}px`,
            ].join(';');
            btn.appendChild(r);
            setTimeout(() => r && r.remove(), 620);
        });
    }

    function attachTilt(el) {
        if (el._fqTilt) return;
        el._fqTilt = true;
        el.addEventListener('mousemove', function(e) {
            const rect = el.getBoundingClientRect();
            const cx = rect.left + rect.width / 2;
            const cy = rect.top  + rect.height / 2;
            const dx = (e.clientX - cx) / (rect.width  / 2);
            const dy = (e.clientY - cy) / (rect.height / 2);
            el.style.transform = `perspective(700px) rotateY(${dx * 4}deg) rotateX(${-dy * 4}deg) scale(1.012)`;
            el.style.transition = 'transform 0.08s linear';
        });
        el.addEventListener('mouseleave', function() {
            el.style.transform = 'perspective(700px) rotateY(0deg) rotateX(0deg) scale(1)';
            el.style.transition = 'transform 0.4s cubic-bezier(0.23,1,0.32,1)';
        });
    }

    function attachFoldIn(el) {
        if (el._fqFold) return;
        el._fqFold = true;
        el.style.animation = 'fold-in 0.45s cubic-bezier(0.23,1,0.32,1) both';
    }

    const mo = new MutationObserver(function() {
        doc.querySelectorAll(
            '[data-testid="stBaseButton-primary"], .stButton > button[kind="primary"]'
        ).forEach(attachRipple);
        doc.querySelectorAll('.stDownloadButton > button').forEach(attachRipple);
        doc.querySelectorAll('[data-testid="stFileUploader"] section').forEach(attachTilt);
        doc.querySelectorAll('[data-testid="stAlert"]').forEach(attachFoldIn);
    });
    mo.observe(doc.body, { childList: true, subtree: true });
})();
</script>
""", height=0)

# ── Hero header ───────────────────────────────────────────────────────────────
_logo_path = BASE_DIR / "logo package" / "PNG" / "primary logo + tagline_reverse.png"
_logo_b64  = base64.b64encode(_logo_path.read_bytes()).decode() if _logo_path.exists() else ""
_logo_img  = (
    f'<img src="data:image/png;base64,{_logo_b64}" '
    f'style="height:38px;width:auto;object-fit:contain;margin-bottom:1rem;display:block;">'
    if _logo_b64 else ""
)

_potx_path   = BASE_DIR / "FulcrumQ_Theme_vF.potx"
_potx_b64    = base64.b64encode(_potx_path.read_bytes()).decode() if _potx_path.exists() else ""
_potx_button = (
    f'<a class="fq-hero-dl" '
    f'href="data:application/vnd.openxmlformats-officedocument.presentationml.template;base64,{_potx_b64}" '
    f'download="FulcrumQ_Template.potx">'
    f'<span class="fq-hero-dl-label">Download template</span>'
    f'<span class="fq-hero-dl-name">FulcrumQ Master</span>'
    f'<span class="fq-hero-dl-ext">.POTX</span>'
    f'</a>'
    if _potx_b64 else ""
)

st.markdown(f"""
<div class="fq-header" style="display:flex;justify-content:space-between;align-items:flex-start;">
    <div>
        {_logo_img}
        <h1>Deck Converter</h1>
        <p>Transform any source deck to FulcrumQ brand in seconds.</p>
        <div class="fq-steps">
            <span class="fq-step-pill">① Master swap</span>
            <span class="fq-step-pill">② Layout remap</span>
            <span class="fq-step-pill">③ Vestige cleanup</span>
            <span class="fq-step-pill">④ Brand style pass</span>
        </div>
    </div>
    {_potx_button}
</div>
<div class="fq-pivot-line"></div>
""", unsafe_allow_html=True)

_tab_labels = ["  Convert  ", "  Color Palette  ", "  Logo Suite  ", "  Icons  ", "  RAG  ", "  Semantic Ingest  ", "  Guided Convert  "]
tab_convert, tab_palette, tab_logos, tab_icons, tab_rag, tab_semantic, tab_guided = st.tabs(_tab_labels)


# ══════════════════════════════════════════════════════════════════════════════
# TAB 1 — CONVERT
# ══════════════════════════════════════════════════════════════════════════════
with tab_convert:
    _cmode = st.radio(
        "_cmode",
        ["Convert", "Pipeline", "Brand Check", "Visual Ingest Test"],
        horizontal=True,
        label_visibility="collapsed",
    )
    st.markdown("<div style='height:6px'></div>", unsafe_allow_html=True)

    # ── Convert mode ──────────────────────────────────────────────────────────
    if _cmode == "Convert":
        master_exists = cd.MASTER_X.exists()

        uploaded = st.file_uploader(
            "Upload source PPTX",
            type=["pptx"],
            accept_multiple_files=True,
            label_visibility="collapsed",
        )

        if uploaded:
            st.markdown("<br>", unsafe_allow_html=True)

            use_vision = False

            col_btn, _ = st.columns([1, 4])
            with col_btn:
                run = st.button(
                    f"▶  Convert  ({len(uploaded)} file{'s' if len(uploaded) != 1 else ''})",
                    type="primary",
                    disabled=not master_exists,
                    use_container_width=True,
                )

            if run:
                with tempfile.TemporaryDirectory() as tmp:
                    tmp_path = Path(tmp)
                    results: list[tuple] = []

                    for uf in uploaded:
                        src_bytes = uf.getvalue()
                        src_path  = tmp_path / uf.name
                        src_path.write_bytes(src_bytes)

                        audit = _audit_source(src_bytes)

                        log_buf = io.StringIO()
                        with st.spinner(f"Converting {uf.name}…"):
                            with contextlib.redirect_stdout(log_buf):
                                try:
                                    out_path = cd.convert(src_path)
                                except Exception as exc:
                                    results.append((uf.name, None, str(exc), audit, [], None))
                                    continue

                        with st.spinner(f"Rendering preview…"):
                            thumbs = _render_slide_thumbs(out_path)

                        # AI vision classification (runs on source file)
                        vis_results = None
                        if use_vision:
                            with st.spinner(f"AI classifying {uf.name}…"):
                                try:
                                    vis_results = classify_deck_vision(src_path)
                                except Exception as _ve:
                                    st.warning(f"Vision classifier failed: {_ve}")

                        results.append((uf.name, out_path, None, audit, thumbs, vis_results))

                        logs = log_buf.getvalue()
                        with st.expander(f"Log — {uf.name}", expanded=False):
                            st.markdown(
                                f'<div class="log-box">{logs}</div>',
                                unsafe_allow_html=True,
                            )

                    n_ok  = sum(1 for _, o, *_ in results if o and o.exists())
                    n_err = len(results) - n_ok
                    if n_ok:
                        st.success(f"✓  {n_ok}/{len(results)} converted")
                    if n_err:
                        st.error(f"✗  {n_err} failed")

                    # Badge colours per slide type
                    _VIS_BADGE = {
                        "cover":   ("#765FFF", "#fff"),
                        "divider": ("#281A42", "#fff"),
                        "ending":  ("#00C27A", "#fff"),
                        "content": ("#E9E4FF", "#281A42"),
                    }

                    for orig_name, out_path, err, audit, thumbs, vis_results in results:
                        if err or not (out_path and out_path.exists()):
                            st.error(f"{orig_name}: {err}")
                            continue

                        st.download_button(
                            label=f"⬇  {out_path.name}",
                            data=out_path.read_bytes(),
                            file_name=out_path.name,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            key=f"dl_{out_path.name}",
                        )

                        # ── Slide preview ──────────────────────────────────
                        if thumbs:
                            st.markdown(
                                '<div class="section-label" style="margin-top:1rem;">Slide preview</div>',
                                unsafe_allow_html=True,
                            )
                            # Build vision lookup {slide_num: result}
                            _vis_map = {r["slide"]: r for r in (vis_results or [])}
                            _THUMB_COLS = 6
                            for _ri in range(0, len(thumbs), _THUMB_COLS):
                                _row = thumbs[_ri: _ri + _THUMB_COLS]
                                _cells = ""
                                for _ti, _tb in enumerate(_row, _ri + 1):
                                    _tb64 = base64.b64encode(_tb).decode()
                                    _vr   = _vis_map.get(_ti)
                                    if _vr:
                                        _bg, _fg = _VIS_BADGE.get(_vr["type"], ("#ccc", "#000"))
                                        _flag    = " ⚑" if _vr["confidence"] < _VIS_THRESHOLD else ""
                                        _badge   = (
                                            f'<div style="font-size:9px;font-weight:600;'
                                            f'background:{_bg};color:{_fg};border-radius:3px;'
                                            f'padding:1px 5px;margin-top:3px;text-align:center;">'
                                            f'{_vr["type"].upper()}{_flag}</div>'
                                        )
                                    else:
                                        _badge = ""
                                    _cells += (
                                        f'<div class="slide-thumb">'
                                        f'<img src="data:image/png;base64,{_tb64}">'
                                        f'<div class="slide-thumb-num">{_ti}</div>'
                                        f'{_badge}'
                                        f'</div>'
                                    )
                                st.markdown(
                                    f'<div style="display:grid;grid-template-columns:repeat({_THUMB_COLS},1fr);'
                                    f'gap:8px;margin-bottom:8px;">{_cells}</div>',
                                    unsafe_allow_html=True,
                                )

                        # ── Vision classification detail ────────────────────
                        if vis_results:
                            with st.expander("AI classification detail", expanded=False):
                                _low = [r for r in vis_results if r["confidence"] < _VIS_THRESHOLD]
                                if _low:
                                    st.warning(f"{len(_low)} slide(s) flagged for review (confidence < {_VIS_THRESHOLD})")
                                _rows = "".join(
                                    f"<tr>"
                                    f"<td>{r['slide']}</td>"
                                    f"<td><b>{r['type']}</b></td>"
                                    f"<td>{r['confidence']:.2f}</td>"
                                    f"<td>{'↺' if r.get('changed') else ''}</td>"
                                    f"<td style='color:#888;font-size:12px'>{r['reasoning'][:80]}</td>"
                                    f"</tr>"
                                    for r in vis_results
                                )
                                st.markdown(
                                    f"<table style='width:100%;font-size:13px;border-collapse:collapse'>"
                                    f"<thead><tr><th>#</th><th>Type</th><th>Conf</th><th></th><th>Reasoning</th></tr></thead>"
                                    f"<tbody>{_rows}</tbody></table>",
                                    unsafe_allow_html=True,
                                )

                        # ── Unmapped report ────────────────────────────────
                        uc  = audit["colors"]
                        uf2 = audit["fonts"]
                        if uc or uf2:
                            with st.expander(
                                f"Unmapped items — {len(uc)} color(s), {len(uf2)} font(s)",
                                expanded=True,
                            ):
                                col_c, col_f = st.columns(2)
                                with col_c:
                                    st.markdown(
                                        '<div class="section-label">Colors not mapped</div>',
                                        unsafe_allow_html=True,
                                    )
                                    if uc:
                                        rows = "".join(
                                            f'<tr><td><span class="audit-dot" style="background:#{h};"></span>#{h}</td>'
                                            f'<td style="color:rgba(255,255,255,0.35);">×{cnt}</td></tr>'
                                            for h, cnt in sorted(uc.items(), key=lambda x: -x[1])
                                        )
                                        st.markdown(
                                            f'<table class="audit-table"><thead><tr><th>Hex</th><th>Uses</th></tr>'
                                            f'</thead><tbody>{rows}</tbody></table>',
                                            unsafe_allow_html=True,
                                        )
                                    else:
                                        st.markdown(
                                            '<span style="font-size:11px;color:var(--fq-green);">All colors mapped ✓</span>',
                                            unsafe_allow_html=True,
                                        )
                                with col_f:
                                    st.markdown(
                                        '<div class="section-label">Fonts not mapped</div>',
                                        unsafe_allow_html=True,
                                    )
                                    if uf2:
                                        rows = "".join(
                                            f'<tr><td style="font-family:\'{f}\',sans-serif;">{f}</td>'
                                            f'<td style="color:rgba(255,255,255,0.35);">×{cnt}</td></tr>'
                                            for f, cnt in sorted(uf2.items(), key=lambda x: -x[1])
                                        )
                                        st.markdown(
                                            f'<table class="audit-table"><thead><tr><th>Typeface</th><th>Uses</th></tr>'
                                            f'</thead><tbody>{rows}</tbody></table>',
                                            unsafe_allow_html=True,
                                        )
                                    else:
                                        st.markdown(
                                            '<span style="font-size:11px;color:var(--fq-green);">All fonts mapped ✓</span>',
                                            unsafe_allow_html=True,
                                        )

    # ── Pipeline mode ─────────────────────────────────────────────────────────
    elif _cmode == "Pipeline":
        _master_ok = _PIPELINE_AVAILABLE and _PIPELINE_MASTER and _PIPELINE_MASTER.exists()

        if not _PIPELINE_AVAILABLE:
            st.error("Pipeline modules not found. Ensure classifier.py, schema_engine.py, renderer.py, qa_validator.py are in the project directory.")
        elif not _master_ok:
            st.error(f"Master PPTX not found: {_PIPELINE_MASTER}")
        else:
            st.markdown(
                '<div class="section-label" style="margin-bottom:0.5rem;">'
                '4-layer pipeline — Classify → Schema → Render → QA</div>',
                unsafe_allow_html=True,
            )

            pl_uploaded = st.file_uploader(
                "Upload source PPTX",
                type=["pptx"],
                accept_multiple_files=False,
                label_visibility="collapsed",
                key="pl_upload",
            )

            if pl_uploaded:
                col_pl, _ = st.columns([1, 4])
                with col_pl:
                    pl_run = st.button(
                        "▶  Run Pipeline",
                        type="primary",
                        use_container_width=True,
                        key="pl_run",
                    )

                if pl_run:
                    with tempfile.TemporaryDirectory() as _tmp:
                        _tmp_path = Path(_tmp)
                        _src = _tmp_path / pl_uploaded.name
                        _src.write_bytes(pl_uploaded.getvalue())
                        _out_pptx = _tmp_path / (pl_uploaded.name.replace(".pptx", "_FQ.pptx"))
                        _cls_json = _tmp_path / "classifier_output.json"
                        _dec_json = _tmp_path / "schema_decisions.json"
                        _ren_json = _tmp_path / "render_log.json"
                        _qa_json  = _tmp_path / "qa_report.json"

                        # Layer 1 — Classify
                        with st.spinner("Layer 1 — Classifying slides…"):
                            _cls_results = _pl_classify(_src, output_path=_cls_json)

                        # Layer 2 — Schema
                        with st.spinner("Layer 2 — Selecting layouts…"):
                            _decisions = _pl_schema(_cls_results, _PIPELINE_RULES, output_path=_dec_json)

                        # Layer 3 — Render
                        with st.spinner("Layer 3 — Rendering…"):
                            _ren_log = _pl_render(_src, _decisions, _PIPELINE_MASTER, _out_pptx, output_log=_ren_json)

                        # Layer 4 — QA
                        with st.spinner("Layer 4 — QA validation…"):
                            _qa_report = _pl_qa(_out_pptx, _ren_log, _cls_results, output_path=_qa_json)

                        # ── Render thumbnails for before/after ─────────────
                        with st.spinner("Rendering preview…"):
                            _before_thumbs = _render_slide_thumbs(_src)
                            _after_thumbs  = _render_slide_thumbs(_out_pptx) if _out_pptx.exists() else []

                        _qs = _qa_report["summary"]
                        _total = _qs["total"]
                        _pass  = _qs["pass"]
                        _rev   = _qs["review"]
                        _fail  = _qs["fail"]

                        # ── Status bar ────────────────────────────────────
                        st.markdown(
                            f'<div style="display:flex;gap:10px;margin:0.75rem 0 1rem;flex-wrap:wrap;">'
                            f'<span class="master-badge master-ok">✓ {_pass} pass</span>'
                            f'<span class="master-badge" style="background:rgba(255,181,71,0.12);border:1px solid rgba(255,181,71,0.4);color:#FFB547;">△ {_rev} review</span>'
                            f'<span class="master-badge master-err">✗ {_fail} fail</span>'
                            f'<span class="master-badge" style="background:rgba(255,255,255,0.04);border:1px solid rgba(255,255,255,0.1);color:rgba(255,255,255,0.4);">{_total} slides</span>'
                            f'</div>',
                            unsafe_allow_html=True,
                        )

                        # ── Download ──────────────────────────────────────
                        if _out_pptx.exists():
                            st.download_button(
                                label=f"⬇  {_out_pptx.name}",
                                data=_out_pptx.read_bytes(),
                                file_name=_out_pptx.name,
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                key="pl_dl",
                            )

                        # ── Before / After thumbnails ─────────────────────
                        if _before_thumbs or _after_thumbs:
                            _qa_by_id = {s["slide_id"]: s for s in _qa_report["slides"]}
                            _cls_by_id = {r["slide_id"]: r for r in _cls_results}
                            _dec_by_id = {d["slide_id"]: d for d in _decisions}

                            _STATUS_COLOR_PILL = {
                                "pass":   ("#00C27A", "#fff"),
                                "review": ("#FFB547", "#1d1d1d"),
                                "fail":   ("#FF2E88", "#fff"),
                            }
                            _TYPE_COLOR_PILL = {
                                "cover":   ("#765FFF", "#fff"),
                                "divider": ("#281A42", "#fff"),
                                "ending":  ("#00C27A", "#fff"),
                                "content": ("#E9E4FF", "#281A42"),
                            }

                            st.markdown(
                                '<div class="section-label" style="margin:1rem 0 0.5rem;">Before / After</div>',
                                unsafe_allow_html=True,
                            )

                            _COLS = 4
                            _n_slides = max(len(_before_thumbs), len(_after_thumbs))
                            for _row_start in range(0, _n_slides, _COLS):
                                _cells_html = ""
                                for _si in range(_row_start, min(_row_start + _COLS, _n_slides)):
                                    _slide_num = _si + 1
                                    _qa_s   = _qa_by_id.get(_slide_num, {})
                                    _cls_s  = _cls_by_id.get(_slide_num, {})
                                    _dec_s  = _dec_by_id.get(_slide_num, {})
                                    _status = _qa_s.get("status", "?")
                                    _stype  = _cls_s.get("type", "?")
                                    _layout = _dec_s.get("selected_layout", "?")
                                    _conf   = _cls_s.get("confidence", 0)
                                    _issues = _qa_s.get("issues", [])

                                    _sc, _sf = _STATUS_COLOR_PILL.get(_status, ("#888", "#fff"))
                                    _tc, _tf = _TYPE_COLOR_PILL.get(_stype, ("#888", "#fff"))
                                    _flag = " ⚑" if _conf < 0.75 else ""

                                    _b_src = base64.b64encode(_before_thumbs[_si]).decode() if _si < len(_before_thumbs) else ""
                                    _a_src = base64.b64encode(_after_thumbs[_si]).decode()  if _si < len(_after_thumbs)  else ""
                                    _b_img = f'<img src="data:image/png;base64,{_b_src}" style="width:100%;display:block;">' if _b_src else '<div style="width:100%;padding-top:56.25%;background:#111;"></div>'
                                    _a_img = f'<img src="data:image/png;base64,{_a_src}" style="width:100%;display:block;">' if _a_src else '<div style="width:100%;padding-top:56.25%;background:#111;"></div>'

                                    _issue_txt = ""
                                    if _issues:
                                        _issue_txt = f'<div style="font-size:8px;color:rgba(255,255,255,0.4);margin-top:2px;line-height:1.3;">{"; ".join(_issues[:2])}{"…" if len(_issues)>2 else ""}</div>'

                                    _cells_html += f"""
                                    <div style="display:flex;flex-direction:column;gap:2px;">
                                      <div style="font-size:8px;color:rgba(255,255,255,0.3);text-align:center;">#{_slide_num} before</div>
                                      <div class="slide-thumb">{_b_img}</div>
                                      <div style="font-size:8px;color:rgba(255,255,255,0.3);text-align:center;">after</div>
                                      <div class="slide-thumb">{_a_img}</div>
                                      <div style="display:flex;gap:3px;flex-wrap:wrap;margin-top:2px;">
                                        <span style="font-size:8px;font-weight:600;background:{_tc};color:{_tf};border-radius:3px;padding:1px 4px;">{_stype.upper()}{_flag}</span>
                                        <span style="font-size:8px;background:{_sc};color:{_sf};border-radius:3px;padding:1px 4px;">{_status.upper()}</span>
                                      </div>
                                      <div style="font-size:8px;color:rgba(255,255,255,0.35);line-height:1.2;">{_layout}</div>
                                      {_issue_txt}
                                    </div>"""

                                st.markdown(
                                    f'<div style="display:grid;grid-template-columns:repeat({_COLS},1fr);gap:10px;margin-bottom:10px;">{_cells_html}</div>',
                                    unsafe_allow_html=True,
                                )

                        # ── JSON artifact expanders ────────────────────────
                        with st.expander("Classifier output", expanded=False):
                            st.json(_cls_results)

                        with st.expander("Schema decisions", expanded=False):
                            st.json(_decisions)

                        with st.expander("QA report", expanded=False):
                            st.json(_qa_report)

                        with st.expander("Render log", expanded=False):
                            st.json(_ren_log)

    # ── Brand Check mode ──────────────────────────────────────────────────────
    elif _cmode == "Brand Check":
        st.markdown(
            '<div class="section-label" style="margin-bottom:0.5rem;">'
            'Upload a deck to flag off-brand colors, fonts, and CEO Works references per slide</div>',
            unsafe_allow_html=True,
        )
        bc_uploaded = st.file_uploader(
            "Upload PPTX for brand check",
            type=["pptx"],
            accept_multiple_files=False,
            label_visibility="collapsed",
            key="bc_upload",
        )
        if bc_uploaded:
            col_bc, _ = st.columns([1, 4])
            with col_bc:
                bc_run = st.button(
                    "Check Brand Compliance",
                    type="primary",
                    use_container_width=True,
                    key="bc_run",
                )
            if bc_run:
                with st.spinner("Scanning slides…"):
                    bc_issues = _brand_check(bc_uploaded.getvalue())

                with zipfile.ZipFile(io.BytesIO(bc_uploaded.getvalue())) as _z:
                    total_slides = sum(
                        1 for k in _z.namelist()
                        if re.match(r"ppt/slides/slide\d+\.xml$", k)
                    )

                clean = total_slides - len(bc_issues)
                st.markdown(
                    f'<div style="display:flex;gap:12px;margin:0.8rem 0;">'
                    f'<span class="master-badge master-ok">✓ {clean} slide{"s" if clean != 1 else ""} clean</span>'
                    f'<span class="master-badge master-err">✗ {len(bc_issues)} slide{"s" if len(bc_issues) != 1 else ""} flagged</span>'
                    f'</div>',
                    unsafe_allow_html=True,
                )

                if not bc_issues:
                    st.success("All slides pass brand check.")
                else:
                    for issue in bc_issues:
                        snum = issue["slide"]
                        n_c  = len(issue["colors"])
                        n_f  = len(issue["fonts"])
                        n_n  = issue["names"]
                        parts = []
                        if n_c: parts.append(f"{n_c} off-brand color{'s' if n_c != 1 else ''}")
                        if n_f: parts.append(f"{n_f} off-brand font{'s' if n_f != 1 else ''}")
                        if n_n: parts.append(f"{n_n} 'CEO Works' ref{'s' if n_n != 1 else ''}")
                        label = f"Slide {snum} — {', '.join(parts)}"

                        with st.expander(label, expanded=False):
                            col_c2, col_f2, col_n2 = st.columns(3)
                            with col_c2:
                                st.markdown(
                                    '<div class="section-label">Off-brand colors</div>',
                                    unsafe_allow_html=True,
                                )
                                if issue["colors"]:
                                    rows_c = "".join(
                                        f'<tr><td><span class="audit-dot" style="background:#{h};"></span>#{h}</td>'
                                        f'<td style="color:rgba(255,255,255,0.35);">×{cnt}</td></tr>'
                                        for h, cnt in sorted(issue["colors"].items(), key=lambda x: -x[1])
                                    )
                                    st.markdown(
                                        f'<table class="audit-table"><thead><tr><th>Hex</th><th>Uses</th></tr>'
                                        f'</thead><tbody>{rows_c}</tbody></table>',
                                        unsafe_allow_html=True,
                                    )
                                else:
                                    st.markdown(
                                        '<span style="font-size:11px;color:var(--fq-green);">None ✓</span>',
                                        unsafe_allow_html=True,
                                    )
                            with col_f2:
                                st.markdown(
                                    '<div class="section-label">Off-brand fonts</div>',
                                    unsafe_allow_html=True,
                                )
                                if issue["fonts"]:
                                    rows_f = "".join(
                                        f'<tr><td style="font-family:\'{f}\',sans-serif;">{f}</td>'
                                        f'<td style="color:rgba(255,255,255,0.35);">×{cnt}</td></tr>'
                                        for f, cnt in sorted(issue["fonts"].items(), key=lambda x: -x[1])
                                    )
                                    st.markdown(
                                        f'<table class="audit-table"><thead><tr><th>Typeface</th><th>Uses</th></tr>'
                                        f'</thead><tbody>{rows_f}</tbody></table>',
                                        unsafe_allow_html=True,
                                    )
                                else:
                                    st.markdown(
                                        '<span style="font-size:11px;color:var(--fq-green);">None ✓</span>',
                                        unsafe_allow_html=True,
                                    )
                            with col_n2:
                                st.markdown(
                                    '<div class="section-label">Brand name issues</div>',
                                    unsafe_allow_html=True,
                                )
                                if n_n:
                                    st.markdown(
                                        f'<span style="font-size:11px;color:var(--fq-magenta);">'
                                        f'{n_n}x "CEO Works" — run Convert to fix</span>',
                                        unsafe_allow_html=True,
                                    )
                                else:
                                    st.markdown(
                                        '<span style="font-size:11px;color:var(--fq-green);">None ✓</span>',
                                        unsafe_allow_html=True,
                                    )

    # ── Visual Ingest Test mode ────────────────────────────────────────────────
    elif _cmode == "Visual Ingest Test":

        _VIT_REGIONS_PROMPT_TMPL = (
            "You are a presentation slide designer. You will receive a rasterized image "
            "of a slide as a file attachment and its cleaned XML structure.\n\n"
            "Your job is to extract a structured region specification for the slide. "
            "Do not render HTML in this step.\n\n"
            "═══════════════════════════════════════\n"
            "FULCRUMQ BRAND GUIDELINES — APPLY EXACTLY\n\n"
            "Canvas: 1280×720px (16:9), fixed size, no scrolling\n\n"
            "COLORS — use only these values\n"
            "- Pivot Purple: #765FFF (primary brand, 15–20% of surface area)\n"
            "- Vector Dark Purple: #281A42 (dark backgrounds — dividers, section slides)\n"
            "- Guiding Grey: #1D1D1D (primary text on light backgrounds)\n"
            "- Grey 1: #3B3B3B\n"
            "- Grey 2: #585858\n"
            "- Grey 3: #767676\n"
            "- Grey Mid: #7A828D\n"
            "- Light Grey: #B2BBCA\n"
            "- Soft Grey: #D0D7DF\n"
            "- Lightest Grey: #F2F2F2\n"
            "- White: #FFFFFF\n"
            "- Anchor Green: #00C27A (accent only, max 3–5% surface)\n"
            "- Momentum Amber: #FFB547 (accent only)\n"
            "- Talent Teal: #60BDBC (accent only)\n"
            "- Signal Magenta: #FF2E88 (accent only)\n"
            "- Resistance Red: #EA3323 (RAG only)\n"
            "- Light Lavender: #E9E4FF\n"
            "- Tint 1: #917FFF / Tint 2: #AD9FFF / Tint 3: #C8BFFF\n\n"
            "COLOR USAGE RULES\n"
            "- Light slides: white/lightest grey background, Guiding Grey text\n"
            "- Dark slides (dividers, section breaks): #281A42 background, white text\n"
            "- Pivot Purple used for titles on light slides, key callout numbers, accent elements\n"
            "- Accent colors (green, amber, teal, magenta) sparingly for data callouts only\n"
            "- Resistance Red (#EA3323) is semantic-only: use it only for negative/at-risk/low-score RAG or status signals\n"
            "- Never use red for headers, decorative bars, section treatments, dividers, generic callouts, or visual emphasis\n"
            "- NEVER use colors not in the palette above\n\n"
            "TYPOGRAPHY — use exact CSS font-family values as shown\n"
            "- Slide title: font-family: 'Segoe UI', sans-serif; font-weight: bold; font-size: 28px; color #1D1D1D on light / #FFFFFF on dark\n"
            "- Subtitle: font-family: 'Segoe UI', sans-serif; font-weight: normal; font-size: 20px; same color as title\n"
            "- Body bullets (level 1): font-family: Arial, sans-serif; font-size: 20px; color #3B3B3B\n"
            "- Body bullets (level 2): font-family: Arial, sans-serif; font-size: 18px; color #3B3B3B\n"
            "- Body bullets (level 3+): font-family: Arial, sans-serif; font-size: 16px; color #3B3B3B\n"
            "- Section labels / kickers: font-family: Arial, sans-serif; font-size: 12px; color #765FFF; text-transform: uppercase; letter-spacing: 0.08em\n"
            "- Table headers: font-family: Arial, sans-serif; font-weight: bold; font-size: 12px; color #FFFFFF; background #281A42\n"
            "- Table body: font-family: Arial, sans-serif; font-size: 12px; alternating #F2F2F2 / #FFFFFF rows\n"
            "- Stat callouts (large numbers): font-family: 'Segoe UI', sans-serif; font-weight: bold; font-size: 48–64px; color #765FFF or a non-red accent\n"
            "Apply these as inline styles or a <style> block — do not rely on browser defaults.\n\n"
            "BULLET STYLE\n"
            "- Bullet glyph: small filled triangle ▲ in Pivot Purple (#765FFF)\n"
            "- Left-align all body text\n"
            "- Indent nested bullets proportionally\n\n"
            "LAYOUT RULES BY SLIDE TYPE\n"
            "- Cover (light): white/lightest grey bg, large bold title top-left, subtitle below, FulcrumQ logo bottom-left\n"
            "- Cover (dark): #281A42 bg, white title, white subtitle\n"
            "- Divider/Section: #281A42 bg, large bold white title bottom-left, lighter subtitle below it, thin purple line bottom-right with small triangle\n"
            "- Content (light): white bg, bold dark title top-left, subtitle below title, content in main body area\n"
            "- Ending: Pivot Purple (#765FFF) bg, large white bold title, white subtitle\n\n"
            "FOOTER (all content slides)\n"
            "- Bottom-left: FULCRUMQ wordmark in black (light slides) or white (dark slides)\n"
            "- Bottom-center: '© 2026 | All Rights Reserved' in Grey Mid, 8px\n"
            "- Bottom-center: slide number in Grey Mid, 8px\n"
            "- Bottom-right: thin horizontal line + small filled triangle in Pivot Purple\n\n"
            "TABLES\n"
            "- Header row: #281A42 fill, white bold Arial 12px text\n"
            "- Body rows: alternating #F2F2F2 and #FFFFFF\n"
            "- Left-align all body cell text\n"
            "- All borders: 1px #D0D7DF\n\n"
            "═══════════════════════════════════════\n\n"
            "LAYOUT COMPOSITION RULES\n"
            "Before writing any HTML, mentally divide the slide into zones based on what you see in the XML and image. Every content region must have a defined zone. Preserve the original slide's composition, spacing relationships, and relative footprint before you optimize anything. Apply these rules:\n"
            "- Fidelity comes first. Recreate the original layout as closely as possible before making it prettier or cleaner\n"
            "- Use the XML x/y/cx/cy geometry to preserve the original region positions and sizes\n"
            "- You may use absolute positioning for macro regions when needed to match the source layout more faithfully\n"
            "- Use CSS Grid or flexbox inside a region once that region's overall footprint is correct\n"
            "- All content must fit within 1280×720px without overflow or scrolling — scale font sizes down if needed to make everything fit\n"
            "- Numbered sections (1, 2, 3 etc) must all be visible and consistently positioned — if you see 3 numbered items in the source, all 3 must appear in the output\n"
            "- When a stat callout ($236M etc) sits next to a diagram, preserve that same row and relative scale — do not stack or drift it away from its partner region\n"
            "- Do not rebalance whitespace aggressively. If the source is dense, the output should also be dense\n"
            "- Preserve banner bars, footer strips, framed screenshots, and wide horizontal bands at their original visual width and height relationships\n"
            "- The overall composition should feel like the same slide, not a redesign. If the source has left content + center table + right screenshot, the HTML must preserve that 3-region split explicitly\n\n"
            "RENDERING RULES\n"
            "- Use only visible content from the image — do not invent content\n"
            "- Identify the correct semantic role of each element from VISUAL SIZE AND POSITION, not XML tag\n"
            "- The title is the largest or most visually prominent text on the slide, almost always in the top portion. Even if the XML tags it as a body placeholder, if it is visually the biggest text and positioned at the top, it is the title. Render it as an <h1> with font-family: 'Segoe UI', sans-serif; font-weight: bold; minimum font-size: 28px.\n"
            "- If the slide contains a screenshot of a table or chart, extract the data from it and render as a native HTML table or chart\n"
            "- If the slide contains an embedded screenshot, screenshot-like panel, or image-heavy region, preserve it as a visually faithful framed region with the same footprint and hierarchy. Do not flatten it into a simplified semantic summary\n"
            "- If a screenshot or diagram cannot be reconstructed exactly from text alone, recreate its frame, headers, key labels, internal sections, and major visual hierarchy so it still looks like the original screenshot region rather than a generic replacement\n"
            "- If the slide contains a diagram, flowchart, or process visual that cannot be reconstructed from text alone, render your best faithful interpretation in native HTML and CSS — use borders, arrows, colored boxes, and nested regions to preserve the same visual logic and density\n"
            "- If the slide is too visually broken to salvage, set data-mode=\"redesign\" on the root div and rebuild from scratch using the correct layout template above\n"
            "- Otherwise set data-mode=\"correct\" on the root div\n"
            "- Preserve all visible text exactly as written — do not paraphrase or summarize\n"
            "- Output only the HTML for the slide content, no explanation, no markdown, no preamble, no <html>/<body> wrapper tags — just the inner div\n\n"
            "CONTENT FIDELITY RULES\n"
            "You must reproduce ALL visible content from the slide, not a summary of it. If section 2 has a data table with numbers like 9250, 7775, 16%, those exact numbers must appear as a native HTML table. If a section contains a screenshot of a complex diagram, render a simplified but structurally faithful version of it in HTML — do not replace it with a single stat or leave it blank. The three content columns must together fill the full 1280×720 canvas from top to bottom — if they look sparse, increase font sizes, increase card padding, or expand diagram elements until the space is used. Loss of content is always wrong. A slide that is too full is better than a slide that is half empty.\n\n"
            "CONTENT COMPLETENESS RULES\n"
            "Before writing any HTML, audit every visible content region in the XML and make a list of everything you can see. For each region, count how many distinct elements are present — a table and a callout stat are two elements, not one. A complex diagram with two sides is one element but must be rendered at full complexity.\n"
            "- If a region contains a table with data, render the full table with all visible rows, columns, and values — never replace it with a single summary stat\n"
            "- If a region contains both a diagram/screenshot AND separate text or numbers below it, both must appear in the output as separate elements\n"
            "- If a diagram has a two-part structure (before/after, today/tomorrow, left/right comparison), render both parts with their labels and internal content — a simplified single-column version is not acceptable\n"
            "- If you can read text inside a screenshot or diagram, that text must appear in your HTML reconstruction\n"
            "- If you can read the numbers in a data element but cannot clearly read the labels beneath them, use the column headers from the table as the labels instead — never invent descriptive labels. For example, if the table had headers 'Current (#)', 'Future (#)', 'Reduction (%)' and values 9250, 7775, 16%, those exact headers and values must appear together as a table, not as separate stat callouts with invented labels.\n"
            "- When content from a screenshot or embedded image is partially readable, reproduce what you can read exactly and omit what you cannot — never substitute invented text for unreadable content. Unreadable regions can be represented as a lightly styled empty box (e.g. a div with a light grey border and no text) rather than fabricated content.\n"
            "- Blank space in your output means you dropped something — go back and find what is missing\n"
            "- The amount of content in your output should feel proportionally similar to the amount of content in the source XML. A content-dense slide should produce a content-dense HTML output\n\n"
            "COPYEDITING RULES\n"
            "Apply these copyediting rules silently to all text in the output — do not mention them, just apply them:\n\n"
            "TITLE CASE (apply to slide titles and section headings):\n"
            "- Capitalize all words except: articles (a, an, the), coordinating conjunctions (and, but, or, nor, for, so, yet), and prepositions (of, in, on, at, to, by, from, with, through, between, about, against, during, without, within, along, across, after, before, behind, below, beneath, beside, beyond, near, since, toward, under, until, upon, via) unless they are the first or last word of the title\n"
            "- Always capitalize the first and last word regardless of what it is\n"
            "- Always capitalize words after a colon or dash\n\n"
            "SENTENCE CASE (apply to subtitles, body text, bullet points, captions):\n"
            "- Capitalize only the first word and proper nouns\n"
            "- Do not capitalize common nouns mid-sentence\n\n"
            "PROTECTED PROPER NOUNS (always use this exact casing everywhere, never alter):\n"
            "- FulcrumQ\n"
            "- Talent to Value\n"
            "- Jobs to Be Done\n"
            "- Value Agenda\n"
            "- Value Hotspots\n\n"
            "PUNCTUATION CLEANUP:\n"
            "- Remove any space before a colon, semicolon, period, comma, question mark, or exclamation mark\n"
            "- Remove double spaces\n"
            "- Use a straight colon with no surrounding spaces on the left: 'Title: Subtitle' not 'Title : Subtitle'\n\n"
            "NEVER ALTER:\n"
            "- Numbers, percentages, dollar amounts\n"
            "- Acronyms (CEO, CFO, CPG, ROI etc)\n"
            "- Text inside data tables or data cells\n"
            "- Brand names other than the protected list above\n"
            "- Any text you are not fully confident you read correctly from the image\n\n"
            "Slide position: slide {n} of {total}\n"
            "Slide type: {slide_type} (cover | content | divider | ending — use \"content\" if unknown)\n\n"
            "Cleaned slide XML:\n"
            "{cleaned_xml}\n\n"
            "Output ---REGIONS--- followed by a valid JSON object with this schema:\n"
            "{\n"
            '  "canvas": {"width": 1280, "height": 720},\n'
            '  "regions": [\n'
            "    {\n"
            '      "id": "region_1",\n'
            '      "kind": "title|text|banner|footer_strip|table|process|screenshot_panel|image_panel|stat|shape",\n'
            '      "x": 0, "y": 0, "w": 100, "h": 50,\n'
            '      "text": "visible text if applicable",\n'
            '      "bg": "#FFFFFF", "fg": "#1D1D1D", "border": "#D0D7DF",\n'
            '      "font_size": 28, "font_weight": 700,\n'
            '      "preserve_as_image": false,\n'
            '      "source_hint": "what source region this came from"\n'
            "    }\n"
            "  ]\n"
            "}\n"
            "- Use pixel coordinates in the 1280×720 canvas\n"
            "- Regions must describe the major slide structure, not every tiny word\n"
            "- Mark screenshot-heavy or image-heavy regions with preserve_as_image=true\n"
            "- The region list should be sufficient to recreate the slide in native PPTX objects and placed images\n\n"
            "NORMALIZATION RULES FOR TARGET MASTER\n"
            "- Do not include legacy source slide-number badges from the top-left; the target master will provide slide numbering\n"
            "- Do not include legacy footer strips or footer text copied from the source; the target master will provide the footer\n"
            "- Strip leading slide numbers from the title text\n"
            "- Title text must be title-cased, not all caps\n"
            "- Numbered step circles used for section labels should use Pivot Purple (#765FFF), not red, unless they encode true negative/status meaning\n\n"
            "Start your response with ---REGIONS---.\n"
            "CRITICAL OUTPUT FORMAT RULES:\n"
            "- After ---REGIONS---, output valid JSON only\n"
            "- Do not use markdown code fences\n"
            "- Do not include HTML in this step\n"
            "- Do not include explanation, audit, or prose before or after the JSON"
        )

        _VIT_HTML_PROMPT_TMPL = (
            "You are a presentation slide designer. You will receive a rasterized image of a slide "
            "as a file attachment plus a structured region specification.\n\n"
            "Your job is to render faithful, brand-compliant HTML from that region specification. "
            "Do not invent new regions. Use the supplied regions as the source of truth for layout.\n\n"
            "Canvas: 1280×720px (16:9), fixed size, no scrolling\n"
            "Use inline styles or a <style> block. Preserve the original composition and density.\n"
            "Screenshot-heavy regions must keep their original footprint and hierarchy.\n"
            "The overall composition should feel like the same slide, not a redesign.\n\n"
            "Region JSON:\n"
            "{region_json}\n\n"
            "First output ---PPTX_LAYOUT--- followed by valid JSON with this schema:\n"
            "{\n"
            '  "layout_name": "Content_Light|Content_Light_Sub|Title_Only_Light|Blank_Light",\n'
            '  "regions": [same schema as input, but refined for final placement]\n'
            "}\n"
            "- Refine spacing, grouping, alignment, and sizing so the slide feels polished and coherent\n"
            "- Keep screenshot-heavy regions marked preserve_as_image=true when appropriate\n"
            "- This refined JSON will be used to build native PPTX objects, so improve composition where helpful\n\n"
            "PPTX LAYOUT RULES\n"
            "- Reserve the bottom footer zone for the target master; do not place content in the bottom 48px of the canvas\n"
            "- Do not include a source footer strip or source slide-number badge in the refined layout\n"
            "- Strip leading slide numbers from the title and return the title in title case\n"
            "- Numbered step circles should be Pivot Purple (#765FFF), not red, unless they are true negative/status indicators\n"
            "- Avoid overlap between regions. Large stat callouts such as $236M in Potential Savings must have their own visible region and must not be clipped or covered\n"
            "- If there is excess blank vertical space in a column, redistribute regions vertically to create breathing room and prevent crowding near the footer\n"
            "- Where a diagram panel has unused vertical space, expand or reposition its internal content to use the space more effectively\n\n"
            "Then output ---HTML--- and raw HTML markup only.\n"
            "Do not wrap the HTML in quotes or JSON.\n"
            "Do not use markdown code fences.\n"
            "The first non-whitespace character after ---HTML--- must be <"
        )

        def _vit_clean_xml(slide_bytes):
            """Return stripped structural XML with geometry so the model can preserve layout."""
            from lxml import etree as _et
            _NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
            _NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
            try:
                _root   = _et.fromstring(slide_bytes)
                _spTree = _root.find(f".//{{{_NS_P}}}spTree")
                if _spTree is None:
                    return "<slide/>"
                _lines = ["<slide>"]
                for _node in _spTree:
                    if not isinstance(_node.tag, str):
                        continue
                    _local = _node.tag.rsplit("}", 1)[-1]
                    if _local not in {"sp", "pic", "graphicFrame"}:
                        continue
                    _attrs = {}
                    _attrs["kind"] = _local
                    _cNvPr = _node.find(f".//{{{_NS_P}}}cNvPr")
                    if _cNvPr is not None:
                        _attrs["id"] = _cNvPr.get("id", "")
                        if _cNvPr.get("name"):
                            _attrs["name"] = _cNvPr.get("name", "")
                    _ph = _node.find(f".//{{{_NS_P}}}ph")
                    if _ph is not None:
                        _attrs["ph"] = _ph.get("type", "body")
                    _xfrm = _node.find(f".//{{{_NS_A}}}xfrm")
                    if _xfrm is not None:
                        _off = _xfrm.find(f"./{{{_NS_A}}}off")
                        _ext = _xfrm.find(f"./{{{_NS_A}}}ext")
                        if _off is not None:
                            _attrs["x"] = _off.get("x", "")
                            _attrs["y"] = _off.get("y", "")
                        if _ext is not None:
                            _attrs["cx"] = _ext.get("cx", "")
                            _attrs["cy"] = _ext.get("cy", "")
                    _texts = [_t.text for _t in _node.iter(f"{{{_NS_A}}}t")
                              if _t.text and _t.text.strip()]
                    _attr_str = " ".join(f'{k}="{v}"' for k, v in _attrs.items())
                    _text_str = " ".join(_texts)
                    if _text_str:
                        _lines.append(f"  <node {_attr_str}>{_text_str}</node>")
                    else:
                        _lines.append(f"  <node {_attr_str}/>")
                _lines.append("</slide>")
                return "\n".join(_lines)
            except Exception:
                return "<slide/>"

        def _vit_rasterize(pptx_path, tmp_dir, dpi=96):
            """PPTX → PNG per slide via slide_vision.rasterize_pptx (LibreOffice)."""
            try:
                from slide_vision import rasterize_pptx as _sv_rasterize
            except ImportError as _ie:
                raise RuntimeError(f"slide_vision not available: {_ie}")
            return _sv_rasterize(pptx_path, tmp_dir)

        _VIT_RELAY_URL  = "https://anyquest-webhook-relay-production-863f.up.railway.app"
        _VIT_AGENT_SLUG = "prompt-executor-qchksn"

        def _vit_extract_regions(raw_text: str) -> dict | None:
            """Recover region JSON reliably from relay/model output."""
            def _strip_fences(text: str) -> str:
                text = text.strip()
                if text.startswith("```"):
                    text = re.sub(r"^```[a-zA-Z0-9_-]*\s*", "", text, count=1)
                    text = re.sub(r"\s*```$", "", text, count=1)
                return text.strip()

            def _decode_wrapped_text(text: str) -> str:
                text = _strip_fences(text)
                # JSON string or quoted payload
                for _ in range(2):
                    _s = text.strip()
                    if not _s:
                        break
                    try:
                        _parsed = json.loads(_s)
                    except Exception:
                        break
                    if isinstance(_parsed, str):
                        text = _parsed
                        continue
                    if isinstance(_parsed, dict):
                        for _key in ("html", "content", "result", "output"):
                            _val = _parsed.get(_key)
                            if isinstance(_val, str):
                                text = _val
                                break
                        else:
                            return _s
                        continue
                    break
                # Common escaped-text fallback from relay/model wrapping
                if "\\n" in text or "\\t" in text or '\\"' in text:
                    try:
                        text = bytes(text, "utf-8").decode("unicode_escape")
                    except Exception:
                        pass
                return _strip_fences(text)

            _body = raw_text or ""
            if "---REGIONS---" in _body:
                _, _, _body = _body.partition("---REGIONS---")
            _body = _decode_wrapped_text(_body)
            try:
                return json.loads(_body)
            except Exception:
                _m = re.search(r"\{[\s\S]*\}", _body)
                if _m:
                    try:
                        return json.loads(_m.group(0))
                    except Exception:
                        return None
            return None

        def _vit_extract_html_and_layout(raw_text: str) -> tuple[str | None, dict | None]:
            """Recover HTML and refined pptx layout reliably from relay/model output."""
            def _strip_fences(text: str) -> str:
                text = text.strip()
                if text.startswith("```"):
                    text = re.sub(r"^```[a-zA-Z0-9_-]*\s*", "", text, count=1)
                    text = re.sub(r"\s*```$", "", text, count=1)
                return text.strip()

            def _decode_wrapped_text(text: str) -> str:
                text = _strip_fences(text)
                for _ in range(2):
                    _s = text.strip()
                    if not _s:
                        break
                    try:
                        _parsed = json.loads(_s)
                    except Exception:
                        break
                    if isinstance(_parsed, str):
                        text = _parsed
                        continue
                    if isinstance(_parsed, dict):
                        for _key in ("html", "content", "result", "output"):
                            _val = _parsed.get(_key)
                            if isinstance(_val, str):
                                text = _val
                                break
                        else:
                            return _s
                        continue
                    break
                if "\\n" in text or "\\t" in text or '\\"' in text:
                    try:
                        text = bytes(text, "utf-8").decode("unicode_escape")
                    except Exception:
                        pass
                return _strip_fences(text)

            _layout = None
            _body = raw_text or ""
            if "---PPTX_LAYOUT---" in _body:
                _, _, _body = _body.partition("---PPTX_LAYOUT---")
                if "---HTML---" in _body:
                    _layout_text, _, _body = _body.partition("---HTML---")
                    _layout_text = _decode_wrapped_text(_layout_text)
                    try:
                        _layout = json.loads(_layout_text)
                    except Exception:
                        _m = re.search(r"\{[\s\S]*\}", _layout_text)
                        if _m:
                            try:
                                _layout = json.loads(_m.group(0))
                            except Exception:
                                _layout = None
            elif "---HTML---" in _body:
                _, _, _body = _body.partition("---HTML---")
            _body = _decode_wrapped_text(_body)
            if "<" in _body:
                return _body[_body.index("<"):].strip(), _layout
            _match = re.search(r"(<div\b[\s\S]*</div>)", _body, re.IGNORECASE)
            if _match:
                return _match.group(1).strip(), _layout
            return None, _layout

        def _vit_regions_to_pptx(slide_png_bytes: bytes, region_spec: dict | None) -> bytes | None:
            """Build a simple single-slide PPTX from structured regions."""
            if not isinstance(region_spec, dict):
                return None
            _regions = region_spec.get("regions")
            if not isinstance(_regions, list) or not _regions:
                return None
            try:
                from pptx import Presentation as _Presentation
                from pptx.enum.text import MSO_ANCHOR as _MSO_ANCHOR, PP_ALIGN as _PP_ALIGN
                from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as _MSO_SHAPE
                from pptx.enum.shapes import PP_PLACEHOLDER as _PP_PLACEHOLDER
                from pptx.util import Emu as _Emu, Pt as _Pt
                from pptx.dml.color import RGBColor as _RGB
                from PIL import Image as _Img
            except Exception:
                return None

            def _noneish(val) -> bool:
                return str(val or "").strip().lower() in {"", "none", "transparent", "null"}

            def _is_redish(val) -> bool:
                _s = str(val or "").strip().lstrip("#").lower()
                return _s in {"ea3323", "ff0000", "c00000", "cc0000", "b4000c", "d12b1f", "ec0000"}

            def _strip_leading_slide_number(text: str) -> str:
                return re.sub(r"^\s*\d+\s*[:.\-\)]?\s*", "", text or "").strip()

            def _is_legacy_title_badge(r: dict) -> bool:
                try:
                    _text = str(r.get("text", "") or "").strip()
                    return (
                        _text.isdigit()
                        and int(_text) < 100
                        and float(r.get("x", 0) or 0) < 90
                        and float(r.get("y", 0) or 0) < 70
                        and float(r.get("w", 0) or 0) <= 70
                        and float(r.get("h", 0) or 0) <= 70
                    )
                except Exception:
                    return False

            def _is_legacy_footer_strip(r: dict) -> bool:
                try:
                    _kind = str(r.get("kind", ""))
                    _text = str(r.get("text", "") or "").lower()
                    return (
                        _kind == "footer_strip"
                        or float(r.get("y", 0) or 0) >= 640
                        or "all rights reserved" in _text
                        or "critical roles" in _text
                    )
                except Exception:
                    return False

            def _is_step_badge(r: dict) -> bool:
                try:
                    _text = str(r.get("text", "") or "").strip()
                    return (
                        _text.isdigit()
                        and int(_text) < 10
                        and float(r.get("w", 0) or 0) <= 80
                        and float(r.get("h", 0) or 0) <= 80
                        and float(r.get("y", 0) or 0) < 300
                    )
                except Exception:
                    return False

            def _region_area(r: dict) -> float:
                try:
                    return max(0.0, float(r.get("w", 0) or 0)) * max(0.0, float(r.get("h", 0) or 0))
                except Exception:
                    return 0.0

            def _normalize_regions(regions: list[dict]) -> list[dict]:
                _out = [dict(r) for r in regions if isinstance(r, dict)]
                _out = [r for r in _out if not _is_legacy_title_badge(r) and not _is_legacy_footer_strip(r)]
                if len(_out) > 1:
                    _canvas_area = 1280 * 720
                    _out = [
                        r for r in _out
                        if not (
                            bool(r.get("preserve_as_image"))
                            and _region_area(r) >= _canvas_area * 0.85
                        )
                    ] or _out
                _footer_y = 640
                for r in _out:
                    _kind = str(r.get("kind", "") or "")
                    _text = str(r.get("text", "") or "")
                    if _kind == "title":
                        r["text"] = cd._to_title_case(_strip_leading_slide_number(_text))
                    if _is_step_badge(r) and _is_redish(r.get("bg")):
                        r["bg"] = "#765FFF"
                        if _noneish(r.get("fg")):
                            r["fg"] = "#FFFFFF"
                    try:
                        _y = float(r.get("y", 0) or 0)
                        _h = float(r.get("h", 0) or 0)
                        if _y + _h > _footer_y:
                            r["y"] = max(0, _footer_y - _h - 8)
                    except Exception:
                        pass
                    if _kind == "stat" or "$" in _text:
                        try:
                            _h = float(r.get("h", 0) or 0)
                            r["h"] = max(_h, 52)
                            _y = float(r.get("y", 0) or 0)
                            if _y + float(r["h"]) > _footer_y:
                                r["y"] = max(0, _footer_y - float(r["h"]) - 8)
                        except Exception:
                            pass
                # Draw image/background regions first, then native regions on top.
                _out.sort(key=lambda r: (0 if bool(r.get("preserve_as_image")) else 1, _region_area(r)))
                return _out

            def _choose_layout_name(regions: list[dict]) -> str:
                _kinds = {str(r.get("kind", "")) for r in regions if isinstance(r, dict)}
                _has_title = "title" in _kinds
                _has_banner = "banner" in _kinds
                _has_footer_strip = "footer_strip" in _kinds
                _image_heavy = sum(1 for r in regions if isinstance(r, dict) and bool(r.get("preserve_as_image"))) >= 1
                if _has_title and (_has_banner or _has_footer_strip or _image_heavy):
                    return "Content_Light"
                if _has_title:
                    return "Title_Only_Light"
                return "Blank_Light"

            def _find_layout(prs, layout_name: str):
                for _layout in prs.slide_layouts:
                    if _layout.name == layout_name:
                        return _layout
                return prs.slide_layouts[0]

            def _set_title_placeholder(slide, text: str) -> bool:
                for _ph in slide.placeholders:
                    try:
                        if _ph.placeholder_format.type == _PP_PLACEHOLDER.TITLE:
                            _ph.text = text
                            return True
                    except Exception:
                        continue
                return False

            def _px_to_emu(val: float | int, axis: str) -> int:
                _canvas = 1280 if axis == "x" else 720
                _size = prs.slide_width if axis == "x" else prs.slide_height
                return int(max(0, float(val or 0)) / _canvas * _size)

            def _hex_rgb(val: str | None, default: str) -> tuple[int, int, int]:
                _h = (val or default).strip().lstrip("#")
                if len(_h) != 6 or not re.fullmatch(r"[0-9A-Fa-f]{6}", _h):
                    _h = default.lstrip("#")
                return int(_h[0:2], 16), int(_h[2:4], 16), int(_h[4:6], 16)

            _template = BASE_DIR / "FulcrumQ Default Template.pptx"
            prs = _Presentation(str(_template)) if _template.exists() else _Presentation()
            prs.slide_width = _Emu(12192000)
            prs.slide_height = _Emu(6858000)
            _layout_name = str(region_spec.get("layout_name") or _choose_layout_name(_regions))
            slide = prs.slides.add_slide(_find_layout(prs, _layout_name))

            _img = _Img.open(io.BytesIO(slide_png_bytes)).convert("RGB")

            for _region in _normalize_regions(_regions):
                x = _px_to_emu(_region.get("x", 0), "x")
                y = _px_to_emu(_region.get("y", 0), "y")
                w = max(_px_to_emu(_region.get("w", 0), "x"), _Emu(45720))
                h = max(_px_to_emu(_region.get("h", 0), "y"), _Emu(45720))
                kind = str(_region.get("kind", "shape"))
                text = str(_region.get("text", "") or "")
                bg = _region.get("bg")
                fg = _region.get("fg")
                border = _region.get("border")
                font_size = _region.get("font_size", 18)
                font_weight = int(_region.get("font_weight", 400) or 400)
                preserve_as_image = bool(_region.get("preserve_as_image"))

                if kind == "title" and text and _set_title_placeholder(slide, text):
                    continue

                if preserve_as_image:
                    _crop = _img.crop((
                        max(0, int(_region.get("x", 0))),
                        max(0, int(_region.get("y", 0))),
                        min(_img.width, int((_region.get("x", 0) or 0) + (_region.get("w", 0) or 0))),
                        min(_img.height, int((_region.get("y", 0) or 0) + (_region.get("h", 0) or 0))),
                    ))
                    _buf = io.BytesIO()
                    _crop.save(_buf, format="PNG")
                    _buf.seek(0)
                    slide.shapes.add_picture(_buf, x, y, width=w, height=h)
                    continue

                _shape = slide.shapes.add_shape(_MSO_SHAPE.RECTANGLE, x, y, w, h)
                if not _noneish(bg):
                    _r, _g, _b = _hex_rgb(bg, "#FFFFFF")
                    _shape.fill.solid()
                    _shape.fill.fore_color.rgb = _RGB(_r, _g, _b)
                else:
                    _shape.fill.background()
                if not _noneish(border):
                    _r, _g, _b = _hex_rgb(border, "#D0D7DF")
                    _shape.line.color.rgb = _RGB(_r, _g, _b)
                else:
                    _shape.line.fill.background()

                if text:
                    _tf = _shape.text_frame
                    _tf.clear()
                    _tf.word_wrap = True
                    _tf.vertical_anchor = _MSO_ANCHOR.MIDDLE
                    _tf.margin_left = _Emu(45720)
                    _tf.margin_right = _Emu(45720)
                    _tf.margin_top = _Emu(22860)
                    _tf.margin_bottom = _Emu(22860)

                    _segments = [text]
                    if kind in {"process", "banner", "footer_strip", "text"} and "|" in text:
                        _segments = [s.strip() for s in text.split("|") if s.strip()]

                    _first = True
                    for _seg in _segments:
                        _p = _tf.paragraphs[0] if _first else _tf.add_paragraph()
                        _first = False
                        _p.alignment = _PP_ALIGN.CENTER if kind in {"title", "stat", "footer_strip", "process"} else _PP_ALIGN.LEFT
                        _run = _p.add_run()
                        _run.text = _seg
                        _font = _run.font
                        _font.name = "Segoe UI" if kind in {"title", "stat"} else "Arial"
                        try:
                            _font.size = _Pt(max(8, min(40, float(font_size))))
                        except Exception:
                            _font.size = _Pt(18)
                        _font.bold = font_weight >= 600 or kind in {"title", "stat"}
                        _r, _g, _b = _hex_rgb(fg, "#1D1D1D")
                        _font.color.rgb = _RGB(_r, _g, _b)

            _out = io.BytesIO()
            prs.save(_out)
            return _out.getvalue()

        def _vit_call_agent(prompt_text, image_bytes, image_name="slide.png"):
            """
            Submit prompt + slide image to AgentRelay via the Job API.
            Mirrors _agent_detect_title in convert_deck.py:
              - POST /submit with agentSlug, Prompt, and image as real file attachment
              - Poll GET /result/{requestId} until completed
            """
            import io as _io
            import json as _json
            import time as _time
            import requests as _req
            from PIL import Image as _Img

            # Compress to JPEG (smaller payload, matches convert_deck.py)
            _img = _Img.open(_io.BytesIO(image_bytes))
            if _img.width > 640:
                _h = int(_img.height * 640 / _img.width)
                _img = _img.resize((640, _h), _Img.LANCZOS)
            _buf = _io.BytesIO()
            _img.convert("RGB").save(_buf, format="JPEG", quality=70)
            _buf.seek(0)

            # Submit
            _resp = _req.post(
                f"{_VIT_RELAY_URL}/submit",
                files=[
                    ("agentId",   (None, "generic-prompt-agent")),
                    ("Prompt",    (None, prompt_text)),
                    ("files",     ("slide.jpg", _buf, "image/jpeg")),
                ],
                timeout=30,
            )
            if not _resp.ok:
                raise Exception(
                    f"AgentRelay /submit returned {_resp.status_code}: {_resp.text[:500]}"
                )
            _result = _resp.json()
            if not _result.get("success"):
                raise Exception(_result.get("error", "Submission failed"))

            _request_id = _result["requestId"]

            # Poll for result
            _timeout_s   = 360
            _deadline    = _time.time() + _timeout_s
            _last_status = None
            while _time.time() < _deadline:
                _poll   = _req.get(f"{_VIT_RELAY_URL}/result/{_request_id}", timeout=10)
                _data   = _poll.json()
                _status = _data.get("status", "unknown")
                if _status != _last_status:
                    print(
                        f"[VIT] poll status={_status!r} t={int(_timeout_s - (_deadline - _time.time()))}s",
                        flush=True,
                    )
                    _last_status = _status
                if _status == "completed":
                    return (
                        _data.get("content")
                        or _data.get("result")
                        or _data.get("output")
                        or ""
                    )
                if _status in ("failed", "not_found"):
                    raise Exception(f"AgentRelay job {_status}: {_data}")
                _time.sleep(3)

            raise TimeoutError(f"AgentRelay: no result after {_timeout_s}s")

        # ── UI ────────────────────────────────────────────────────────────────
        st.caption(
            "Rasterize each slide to PNG, extract structural XML, and render via AgentRelay. "
            "Results stream in slide-by-slide."
        )

        if "vit_results" not in st.session_state:
            st.session_state["vit_results"] = []

        _vit_upload = st.file_uploader(
            "Upload PPTX", type=["pptx"], key="vit_pptx",
        )

        if st.button("Visual Ingest Test", type="primary", key="vit_run",
                     disabled=_vit_upload is None):
            with tempfile.TemporaryDirectory() as _vit_tmp:
                _vit_tmp_p = Path(_vit_tmp)
                _vit_src   = _vit_tmp_p / _vit_upload.name
                _vit_src.write_bytes(_vit_upload.getvalue())
                st.session_state["vit_results"] = []

                _vit_status = st.empty()

                # ── Rasterize all slides
                _vit_status.info("⏳ Rasterizing slides via LibreOffice…")
                try:
                    _vit_pngs = _vit_rasterize(_vit_src, _vit_tmp_p)
                except Exception as _vit_re:
                    st.error(f"Rasterization failed: {_vit_re}")
                    _vit_pngs = []

                if _vit_pngs:
                    # ── Extract cleaned XML for each slide
                    _vit_xmls = []
                    with zipfile.ZipFile(_vit_src) as _vz:
                        _vit_snames = sorted(
                            [_n for _n in _vz.namelist()
                             if re.match(r"ppt/slides/slide\d+\.xml$", _n)],
                            key=lambda _p: int(re.search(r"\d+", _p).group()),
                        )
                        for _sn in _vit_snames:
                            _vit_xmls.append(_vit_clean_xml(_vz.read(_sn)))

                    _vit_n    = min(len(_vit_pngs), len(_vit_xmls))
                    _vit_grid = st.container()

                    for _vi in range(_vit_n):
                        _vit_status.info(f"⏳ Processing slide {_vi + 1} of {_vit_n}…")

                        _v_png_bytes = _vit_pngs[_vi].read_bytes()

                        _v_regions_prompt = (
                            _VIT_REGIONS_PROMPT_TMPL
                            .replace("{n}",            str(_vi + 1))
                            .replace("{total}",        str(_vit_n))
                            .replace("{slide_type}",   "content")
                            .replace("{cleaned_xml}",  _vit_xmls[_vi])
                        )
                        st.caption(
                            f"Slide {_vi + 1}: image {len(_v_png_bytes):,}B "
                            f"prompt {len(_v_regions_prompt):,}c"
                        )

                        try:
                            _v_regions_raw = _vit_call_agent(
                                _v_regions_prompt,
                                _v_png_bytes,
                                image_name=f"slide_{_vi + 1}.png",
                            )
                            _v_regions = _vit_extract_regions(_v_regions_raw)
                            _v_html_raw = ""
                            _v_html = None
                            _v_refined_layout = None
                            if _v_regions:
                                _v_html_prompt = _VIT_HTML_PROMPT_TMPL.replace(
                                    "{region_json}",
                                    json.dumps(_v_regions, indent=2),
                                )
                                _v_html_raw = _vit_call_agent(
                                    _v_html_prompt,
                                    _v_png_bytes,
                                    image_name=f"slide_{_vi + 1}.png",
                                )
                                _v_html, _v_refined_layout = _vit_extract_html_and_layout(_v_html_raw)
                            _v_pptx = _vit_regions_to_pptx(_v_png_bytes, _v_refined_layout or _v_regions)
                            if not _v_html:
                                print(
                                    f"[VIT] Slide {_vi + 1}: could not extract HTML from response",
                                    flush=True,
                                )
                            _v_raw = _v_html_raw or _v_regions_raw
                            _v_err  = None
                        except Exception as _v_ae:
                            _v_html = None
                            _v_regions = None
                            _v_refined_layout = None
                            _v_pptx = None
                            _v_raw = ""
                            _v_err  = (
                                "timeout" if "timeout" in str(_v_ae).lower()
                                else str(_v_ae)
                            )
                        st.session_state["vit_results"].append({
                            "slide_num": _vi + 1,
                            "slide_total": _vit_n,
                            "png_bytes": _v_png_bytes,
                            "html": _v_html,
                            "regions": _v_regions,
                            "refined_layout": _v_refined_layout,
                            "pptx": _v_pptx,
                            "raw": _v_raw,
                            "err": _v_err,
                        })

                    _vit_status.success(f"✓ Done — {_vit_n} slides processed.")

        for _idx, _res in enumerate(st.session_state.get("vit_results", [])):
            with _vit_grid if "_vit_grid" in locals() else st.container():
                with st.container():
                    _slide_num = _res["slide_num"]
                    _slide_total = _res["slide_total"]
                    _v_png_bytes = _res["png_bytes"]
                    _v_html = _res["html"]
                    _v_regions = _res["regions"]
                    _v_refined_layout = _res["refined_layout"]
                    _v_pptx = _res["pptx"]
                    _v_raw = _res["raw"]
                    _v_err = _res["err"]
                    st.markdown(f"---\n#### Slide {_slide_num} of {_slide_total}")
                    _vc1, _vc2 = st.columns(2)
                    with _vc1:
                        st.caption("Original")
                        st.image(_v_png_bytes, use_container_width=True)
                    with _vc2:
                        st.caption("AI Render")
                        if _v_html:
                            _v_html_clean = _v_html
                            _v_shared_style = (
                                '<style>'
                                '*{box-sizing:border-box;}'
                                'body{margin:0;padding:0;overflow:hidden;}'
                                'h1,h2,h3,h4,[class*="title"],[class*="subtitle"],'
                                '[class*="stat"],[class*="callout"]{'
                                "font-family:'Segoe UI',sans-serif!important;}"
                                'p,li,td,th,span,div,[class*="body"],'
                                '[class*="bullet"],[class*="kicker"]{'
                                'font-family:Arial,sans-serif!important;}'
                                '</style>'
                            )
                            _v_frame = (
                                _v_shared_style
                                + '<div style="width:1280px;height:720px;'
                                'transform:scale(0.5);transform-origin:top left;'
                                'overflow:hidden;">'
                                + _v_html_clean
                                + '</div>'
                            )
                            components.html(_v_frame, height=360, scrolling=False)
                            _v_dl_doc = (
                                '<!DOCTYPE html><html lang="en"><head>'
                                '<meta charset="UTF-8">'
                                '<meta name="viewport" content="width=1280">'
                                '<title>Slide ' + str(_slide_num) + '</title>'
                                + _v_shared_style.replace(
                                    'body{margin:0;padding:0;overflow:hidden;}',
                                    'body{margin:0;padding:0;width:1280px;height:720px;overflow:hidden;}'
                                )
                                + '</head><body>'
                                + _v_html_clean
                                + '</body></html>'
                            )
                            st.download_button(
                                label=f"⬇ slide_{_slide_num:02d}.html",
                                data=_v_dl_doc.encode("utf-8"),
                                file_name=f"slide_{_slide_num:02d}.html",
                                mime="text/html",
                                key=f"vit_dl_{_idx}",
                            )
                            if _v_regions:
                                st.download_button(
                                    label=f"⬇ slide_{_slide_num:02d}_regions.json",
                                    data=json.dumps(_v_regions, indent=2).encode("utf-8"),
                                    file_name=f"slide_{_slide_num:02d}_regions.json",
                                    mime="application/json",
                                    key=f"vit_regions_{_idx}",
                                )
                            if _v_refined_layout:
                                st.download_button(
                                    label=f"⬇ slide_{_slide_num:02d}_pptx_layout.json",
                                    data=json.dumps(_v_refined_layout, indent=2).encode("utf-8"),
                                    file_name=f"slide_{_slide_num:02d}_pptx_layout.json",
                                    mime="application/json",
                                    key=f"vit_layout_{_idx}",
                                )
                            if _v_pptx:
                                st.download_button(
                                    label=f"⬇ slide_{_slide_num:02d}.pptx",
                                    data=_v_pptx,
                                    file_name=f"slide_{_slide_num:02d}.pptx",
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                    key=f"vit_pptx_{_idx}",
                                )
                            with st.expander("Raw response", expanded=False):
                                st.code(_v_raw[:12000] or "(empty)", language="text")
                            if _v_regions:
                                with st.expander("Structured regions", expanded=False):
                                    st.json(_v_regions)
                            if _v_refined_layout:
                                with st.expander("Refined PPTX layout", expanded=False):
                                    st.json(_v_refined_layout)
                        elif _v_raw:
                            st.warning("Model returned content, but it was not valid raw HTML.")
                            with st.expander("Raw response", expanded=True):
                                st.code(_v_raw[:12000], language="text")
                        elif _v_err == "timeout":
                            st.warning("⏱ API request timed out for this slide.")
                        else:
                            st.error(f"Agent error: {_v_err}")
                            if _v_raw:
                                with st.expander("Raw response", expanded=False):
                                    st.code(_v_raw[:12000], language="text")


# ══════════════════════════════════════════════════════════════════════════════
# TAB 2 — COLOR PALETTE
# ══════════════════════════════════════════════════════════════════════════════
with tab_palette:

    pal_view = st.radio(
        "palette_view",
        ["Brand Palette", "Conversion Map"],
        horizontal=True,
        label_visibility="collapsed",
    )

    if pal_view == "Brand Palette":

        BRAND_GROUPS = [
            ("Guiding Grey", [
                ("1D1D1D", "Guiding Grey"),
                ("3B3B3B", "Grey 1"),
                ("585858", "Grey 2"),
                ("767676", "Grey 3"),
                ("FFFFFF", "White"),
            ]),
            ("Neutral", [
                ("281A42", "Vector Dark Purple"),
                ("7A828D", "Grey Mid"),
                ("B2BBCA", "Light Grey"),
                ("D0D7DF", "Soft Grey"),
                ("F2F2F2", "Lightest Grey"),
            ]),
            ("Pivot Purple", [
                ("765FFF", "Pivot Purple"),
                ("917FFF", "Tint 1"),
                ("AD9FFF", "Tint 2"),
                ("C8BFFF", "Tint 3"),
                ("E9E4FF", "Light Lavender"),
            ]),
            ("Accents", [
                ("00C27A", "Anchor Green"),
                ("FFB547", "Momentum Amber"),
                ("EA3323", "Resistance Red"),
                ("60BDBC", "Talent Teal"),
                ("FF2E88", "Signal Magenta"),
            ]),
        ]

        def _text_color(hex_val: str) -> str:
            try:
                r = int(hex_val[0:2], 16)
                g = int(hex_val[2:4], 16)
                b = int(hex_val[4:6], 16)
                lum = 0.299 * r + 0.587 * g + 0.114 * b
                return "#1D1D1D" if lum > 140 else "#FFFFFF"
            except Exception:
                return "#FFFFFF"

        st.markdown("""
<style>
.pal-group-label {
    font-family: var(--fq-font-head);
    font-size: 9px; font-weight: 700;
    text-transform: uppercase; letter-spacing: 0.14em;
    color: var(--fq-grey3); margin: 1.4rem 0 0.6rem;
}
.pal-grid {
    display: grid;
    grid-template-columns: repeat(5, 1fr);
    gap: 8px;
}
.pal-tile {
    border-radius: 12px; height: 96px;
    display: flex; flex-direction: column;
    justify-content: flex-end; padding: 8px 10px;
    cursor: crosshair; border: 1px solid rgba(255,255,255,0.07);
    transition: transform 0.18s cubic-bezier(0.23,1,0.32,1),
                box-shadow 0.18s cubic-bezier(0.23,1,0.32,1);
    position: relative; overflow: hidden;
}
.pal-tile:hover {
    transform: translateY(-4px) scale(1.04);
    box-shadow: 0 12px 32px rgba(0,0,0,0.45);
    z-index: 2; border-color: rgba(255,255,255,0.22);
}
.pal-tile-hex {
    font-family: var(--fq-font-mono); font-size: 10px;
    font-weight: 700; letter-spacing: 0.04em; line-height: 1.2;
}
.pal-tile-name {
    font-family: var(--fq-font-body); font-size: 9px;
    opacity: 0.72; line-height: 1.2; margin-top: 1px;
}
</style>
""", unsafe_allow_html=True)

        seen = set()
        for group_label, colors in BRAND_GROUPS:
            st.markdown(
                f'<div class="pal-group-label">{group_label}</div>',
                unsafe_allow_html=True,
            )
            tiles = ""
            for hex_val, name in colors:
                if hex_val in seen:
                    continue
                seen.add(hex_val)
                tc = _text_color(hex_val)
                border = "1px solid rgba(0,0,0,0.15)" if hex_val == "FFFFFF" else "1px solid rgba(255,255,255,0.07)"
                tiles += (
                    f'<div class="pal-tile" style="background:#{hex_val};border:{border};">'
                    f'<span class="pal-tile-hex" style="color:{tc};">#{hex_val}</span>'
                    f'<span class="pal-tile-name" style="color:{tc};">{name}</span>'
                    f'</div>'
                )
            st.markdown(f'<div class="pal-grid">{tiles}</div>', unsafe_allow_html=True)

    else:
        # Conversion Map — group COLOR_REMAP entries by target brand color
        _TARGET_ORDER = [
            ("1D1D1D", "→ Guiding Grey"), ("3B3B3B", "→ Grey 1"),
            ("585858", "→ Grey 2"), ("767676", "→ Grey 3"),
            ("7A828D", "→ Grey Mid"), ("B2BBCA", "→ Light Grey"),
            ("D0D7DF", "→ Soft Grey"), ("FFFFFF", "→ White"),
            ("281A42", "→ Vector Dark Purple"), ("765FFF", "→ Pivot Purple"),
            ("917FFF", "→ Purple Tint 1"), ("AD9FFF", "→ Purple Tint 2"),
            ("C8BFFF", "→ Purple Tint 3"), ("E9E4FF", "→ Shift Lavender"),
            ("00C27A", "→ Anchor Green"), ("60BDBC", "→ Teal"),
            ("FFB547", "→ Amber"), ("FF2E88", "→ Signal Magenta"),
        ]

        def swatch_row(src_hex: str, tgt_hex: str) -> str:
            return (
                f'<div class="swatch-row">'
                f'<span class="swatch" style="background:#{src_hex};"></span>'
                f'<code>#{src_hex}</code>'
                f'<span class="arrow">→</span>'
                f'<span class="swatch" style="background:#{tgt_hex};"></span>'
                f'<code>#{tgt_hex}</code>'
                f'</div>'
            )

        _by_target: dict[str, list[str]] = {t: [] for t, _ in _TARGET_ORDER}
        for src, tgt in cd.COLOR_REMAP.items():
            if tgt in _by_target:
                _by_target[tgt].append(src)

        cols = st.columns(3)
        for _ci, (tgt, label) in enumerate(_TARGET_ORDER):
            srcs = _by_target.get(tgt, [])
            if not srcs:
                continue
            with cols[_ci % 3]:
                st.markdown(f'<div class="section-label">{label}</div>', unsafe_allow_html=True)
                st.markdown("".join(swatch_row(s, tgt) for s in sorted(srcs)), unsafe_allow_html=True)

        st.divider()
        st.markdown('<div class="section-label">Scheme fill overrides</div>', unsafe_allow_html=True)
        scheme_html = "".join(
            f'<div class="scheme-row"><code>{k}</code>'
            f'<span class="arrow">→</span>'
            f'<span class="swatch" style="background:#{v};"></span>'
            f'<code>#{v}</code></div>'
            for k, v in cd.SCHEME_FILL_MAP.items()
        )
        st.markdown(scheme_html, unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# TAB 3 — LOGO SUITE
# ══════════════════════════════════════════════════════════════════════════════
with tab_logos:
    LOGO_DIR = BASE_DIR / "logo package"

    VARIANT_BG = {
        "default":   "#FFFFFF",
        "alternate": "#281A42",
        "reverse":   "#1D1D1D",
    }
    VARIANT_LABEL = {
        "default":   "Default",
        "alternate": "Alternate",
        "reverse":   "Reverse",
    }
    VARIANT_LABEL_COLOR = {
        "default":   "#281A42",
        "alternate": "rgba(255,255,255,0.6)",
        "reverse":   "rgba(255,255,255,0.6)",
    }

    FAMILIES = [
        ("Primary Logo + Tagline", "primary logo + tagline"),
        ("Primary Logo",           "primary logo"),
        ("Monogram",               "monogram"),
        ("Triangle",               "triangle"),
    ]
    VARIANTS = ["default", "alternate", "reverse"]
    DL_FMTS  = [("PNG", "png"), ("SVG", "svg"), ("PDF", "pdf")]

    def _b64(path: Path) -> str:
        return base64.b64encode(path.read_bytes()).decode()

    def _dl_button(path: Path, fmt: str, key: str):
        mime_map = {
            "png": "image/png",
            "svg": "image/svg+xml",
            "pdf": "application/pdf",
        }
        st.download_button(
            label=fmt,
            data=path.read_bytes(),
            file_name=path.name,
            mime=mime_map.get(fmt.lower(), "application/octet-stream"),
            key=key,
            use_container_width=True,
        )

    for family_label, stem in FAMILIES:
        with st.expander(family_label, expanded=(stem == "primary logo + tagline")):
            cols = st.columns(3)
            for col, variant in zip(cols, VARIANTS):
                png_path = LOGO_DIR / "PNG" / f"{stem}_{variant}.png"
                bg  = VARIANT_BG[variant]
                ltc = VARIANT_LABEL_COLOR[variant]

                with col:
                    if png_path.exists():
                        img_b64 = _b64(png_path)
                        st.markdown(
                            f'<div class="logo-tile" style="background:{bg};">'
                            f'<img src="data:image/png;base64,{img_b64}" '
                            f'style="max-width:100%;max-height:120px;object-fit:contain;">'
                            f'<span class="logo-variant-label" style="color:{ltc};">'
                            f'{VARIANT_LABEL[variant]}</span>'
                            f'</div>',
                            unsafe_allow_html=True,
                        )
                    else:
                        st.markdown(
                            f'<div class="logo-tile" style="background:{bg};min-height:120px;">'
                            f'<span style="color:rgba(255,255,255,0.2);font-size:11px;">not found</span>'
                            f'</div>',
                            unsafe_allow_html=True,
                        )

                    st.markdown("<div style='height:6px'></div>", unsafe_allow_html=True)
                    dl_cols = st.columns(len(DL_FMTS))
                    for dc, (fmt_label, fmt_ext) in zip(dl_cols, DL_FMTS):
                        dl_path = LOGO_DIR / fmt_ext.upper() / f"{stem}_{variant}.{fmt_ext}"
                        with dc:
                            if dl_path.exists():
                                _dl_button(dl_path, fmt_label, f"dl_{stem}_{variant}_{fmt_ext}")
                            else:
                                st.markdown(
                                    f'<div style="text-align:center;font-size:9px;'
                                    f'color:rgba(255,255,255,0.18);padding:6px 0;">—</div>',
                                    unsafe_allow_html=True,
                                )


# ══════════════════════════════════════════════════════════════════════════════
# TAB 4 — ICONS
# ══════════════════════════════════════════════════════════════════════════════
with tab_icons:
    import io as _io
    import numpy as _np
    from pptx import Presentation as _Prs
    from pptx.enum.shapes import MSO_SHAPE_TYPE as _MSO
    from PIL import Image as _Image

    _ONEDRIVE_GUIDE = Path(
        "/Users/zervolino/Library/CloudStorage"
        "/OneDrive-SharedLibraries-CEOWorks"
        "/Ming Yang - FulcrumQ Assests"
        "/FulcrumQ PPTX - Master & Usage"
        "/FulcrumQ Slide Guide & Examples.pptx"
    )
    _ICONS_PATH = _ONEDRIVE_GUIDE if _ONEDRIVE_GUIDE.exists() else BASE_DIR / "FQ_PPTX_Theme_vF.pptx"
    _DUO_DARK  = _np.array([0x76, 0x5F, 0xFF], dtype=_np.float32)
    _DUO_LIGHT = _np.array([0xFF, 0xFF, 0xFF], dtype=_np.float32)

    def _apply_duotone(img: "_Image.Image") -> "_Image.Image":
        rgba = _np.array(img.convert("RGBA"), dtype=_np.float32)
        gray = (rgba[:, :, 0] * 0.299 + rgba[:, :, 1] * 0.587 + rgba[:, :, 2] * 0.114) / 255.0
        t = gray[:, :, _np.newaxis]
        rgb_out = (_DUO_DARK * (1 - t) + _DUO_LIGHT * t).clip(0, 255).astype(_np.uint8)
        result = _np.dstack([rgb_out, rgba[:, :, 3].astype(_np.uint8)])
        return _Image.fromarray(result, "RGBA")

    _WIDE_SHAPES = set()   # no wide composites in theme deck
    _EXCLUDED    = set()   # no manual exclusions needed
    _ICON_SIZE   = 120

    def _fit_to_square(img: "_Image.Image", size: int) -> "_Image.Image":
        img.thumbnail((size, size), _Image.LANCZOS)
        canvas = _Image.new("RGBA", (size, size), (0, 0, 0, 0))
        x = (size - img.width) // 2
        y = (size - img.height) // 2
        canvas.paste(img, (x, y))
        return canvas

    def _fit_to_height(img: "_Image.Image", h: int) -> "_Image.Image":
        ratio = h / img.height
        return img.resize((max(1, int(img.width * ratio)), h), _Image.LANCZOS)

    @st.cache_data(show_spinner=False)
    def _load_icons(path_str: str):
        """Return (sections, wide_img).
        sections = [{"title": str, "icons": [PIL.Image]}]
        wide_img  = PIL.Image | None  — rendered separately at the end.
        Text shapes within the slide act as section dividers (sorted by position).
        Blob-hash deduplication prevents identical images appearing twice."""
        import hashlib
        prs = _Prs(path_str)
        seen: set[str] = set()
        sections: list[dict] = []
        cur_title = ""
        cur_icons: list = []
        wide_img = None

        # Only process slides that are part of the Icon Library
        # (skips color/usage/photo slides in the full theme deck)
        icon_slides = []
        for slide in prs.slides:
            for s in slide.shapes:
                try:
                    if "icon" in s.text_frame.text.strip().lower():
                        icon_slides.append(slide)
                        break
                except Exception:
                    pass
        # Fall back to all slides if none matched (e.g. icons_fq.pptx)
        if not icon_slides:
            icon_slides = list(prs.slides)

        for slide in icon_slides:
            for shape in sorted(slide.shapes, key=lambda s: (s.top or 0, s.left or 0)):
                # ── Text shape → start a new section ─────────────────────────
                if shape.has_text_frame and shape.shape_type != _MSO.PICTURE:
                    text = shape.text_frame.text.strip()
                    if len(text) > 1:           # skip page numbers / single chars
                        if cur_icons:
                            sections.append({"title": cur_title, "icons": cur_icons})
                            cur_icons = []
                        cur_title = text
                    continue

                if shape.shape_type != _MSO.PICTURE:
                    continue
                if shape.name in _EXCLUDED:
                    continue

                try:
                    blob = shape.image.blob
                    bh   = hashlib.md5(blob).hexdigest()
                    if bh in seen:
                        continue
                    seen.add(bh)

                    raw = _Image.open(_io.BytesIO(blob)).convert("RGBA")
                    img = _apply_duotone(raw)

                    if shape.name in _WIDE_SHAPES:
                        wide_img = _fit_to_height(img, _ICON_SIZE)
                    else:
                        cur_icons.append(_fit_to_square(img, _ICON_SIZE))
                except Exception:
                    pass

        if cur_icons:
            sections.append({"title": cur_title, "icons": cur_icons})

        return sections, wide_img

    st.markdown("### Icon Library")

    if not _ICONS_PATH.exists():
        st.error(f"File not found: `{_ICONS_PATH}`")
    else:
        _sections, _wide_img = _load_icons(str(_ICONS_PATH))
        _COLS_PER_ROW = 10

        for _sec in _sections:
            if _sec["title"]:
                st.markdown(
                    f'<div class="section-label" style="margin-top:1.4rem;">{_sec["title"]}</div>',
                    unsafe_allow_html=True,
                )
            if _sec["icons"]:
                _cells = ""
                for _img in _sec["icons"]:
                    _buf = _io.BytesIO()
                    _img.save(_buf, format="PNG")
                    _ib64 = base64.b64encode(_buf.getvalue()).decode()
                    _cells += (
                        f'<div class="icon-cell">'
                        f'<img src="data:image/png;base64,{_ib64}" '
                        f'width="{_ICON_SIZE}" height="{_ICON_SIZE}" style="display:block;">'
                        f'</div>'
                    )
                st.markdown(
                    f'<div style="display:grid;grid-template-columns:repeat({_COLS_PER_ROW},1fr);'
                    f'gap:6px;margin-bottom:4px;">{_cells}</div>',
                    unsafe_allow_html=True,
                )

        # ── Wide composite — own full-width row at the end ────────────────────
        if _wide_img is not None:
            st.markdown(
                '<div class="section-label" style="margin-top:1.4rem;">Click process</div>',
                unsafe_allow_html=True,
            )
            _buf = _io.BytesIO()
            _wide_img.save(_buf, format="PNG")
            _wb64 = base64.b64encode(_buf.getvalue()).decode()
            st.markdown(
                f'<div class="icon-cell" style="display:inline-flex;align-items:center;'
                f'padding:10px 14px;height:{_ICON_SIZE}px;">'
                f'<img src="data:image/png;base64,{_wb64}" '
                f'style="height:{_ICON_SIZE - 20}px;width:auto;display:block;">'
                f'</div>',
                unsafe_allow_html=True,
            )

# ══════════════════════════════════════════════════════════════════════════════
# TAB 6 — RAG
# ══════════════════════════════════════════════════════════════════════════════
with tab_rag:

    # ── RAG color scale (red → amber → green using approved brand hex) ─────────
    RAG_SCALE = [
        {"hex": "EA3323", "label": "Critical",   "score": 0},    # Resistance Red
        {"hex": "FFB547", "label": "Caution",     "score": 50},   # Momentum Amber
        {"hex": "00C27A", "label": "Strong",      "score": 100},  # Anchor Green
    ]

    def _rag_color(score: float) -> str:
        """Interpolate hex color for a 0-100 score along the RAG scale."""
        for i in range(len(RAG_SCALE) - 1):
            lo = RAG_SCALE[i]
            hi = RAG_SCALE[i + 1]
            if lo["score"] <= score <= hi["score"]:
                t = (score - lo["score"]) / (hi["score"] - lo["score"])
                def _lerp_chan(a: str, b: str, ti: float, sl: slice) -> int:
                    return round(int(a[sl], 16) * (1 - ti) + int(b[sl], 16) * ti)
                r = _lerp_chan(lo["hex"], hi["hex"], t, slice(0, 2))
                g = _lerp_chan(lo["hex"], hi["hex"], t, slice(2, 4))
                b = _lerp_chan(lo["hex"], hi["hex"], t, slice(4, 6))
                return f"#{r:02X}{g:02X}{b:02X}"
        return f"#{RAG_SCALE[-1]['hex']}"

    def _rag_text(hex_col: str) -> str:
        r, g, b = int(hex_col[1:3], 16), int(hex_col[3:5], 16), int(hex_col[5:7], 16)
        return "#1D1D1D" if (0.299 * r + 0.587 * g + 0.114 * b) > 160 else "#FFFFFF"

    st.markdown("""
<style>
.rag-section-label {
    font-family: var(--fq-font-head);
    font-size: 9px; font-weight: 700; text-transform: uppercase;
    letter-spacing: 0.14em; color: var(--fq-grey3);
    margin: 1.6rem 0 0.7rem;
}
/* ── Color scale strip ── */
.rag-strip { display:flex; border-radius:12px; overflow:hidden; height:72px; margin-bottom:0.5rem; }
.rag-strip-seg {
    flex:1; display:flex; flex-direction:column;
    justify-content:flex-end; padding:7px 10px;
}
.rag-strip-label { font-size:9px; font-weight:700; letter-spacing:0.05em; }
.rag-strip-hex   { font-family:var(--fq-font-mono); font-size:8px; opacity:0.72; margin-top:2px; }
/* ── Gradient bar ── */
.rag-grad-bar {
    height:28px; border-radius:8px;
    background: linear-gradient(to right,
        #EA3323 0%, #FFB547 50%, #00C27A 100%);
    position:relative; margin: 0.4rem 0 1.4rem;
}
.rag-grad-thumb {
    position:absolute; top:-6px; width:16px; height:40px;
    background:#FFFFFF; border-radius:4px;
    box-shadow:0 2px 8px rgba(0,0,0,0.5);
    transform:translateX(-50%);
    border:2px solid rgba(0,0,0,0.25);
    transition: left 0.1s;
}
.rag-grad-labels {
    display:flex; justify-content:space-between;
    font-size:9px; color:rgba(255,255,255,0.45);
    font-family:var(--fq-font-mono); margin-top:2px;
}
/* ── Sample data grid ── */
.rag-table { width:100%; border-collapse:collapse; font-size:12px; }
.rag-table th {
    text-align:left; padding:6px 10px;
    border-bottom:1px solid rgba(255,255,255,0.1);
    font-family:var(--fq-font-head); font-size:9px; font-weight:700;
    text-transform:uppercase; letter-spacing:0.1em;
    color:rgba(255,255,255,0.45);
}
.rag-table td { padding:6px 10px; color:rgba(255,255,255,0.85); }
.rag-table tr:hover td { background:rgba(118,95,255,0.06); }
.rag-badge {
    display:inline-block; border-radius:6px;
    padding:2px 9px; font-weight:700; font-size:11px;
    font-family:var(--fq-font-mono); white-space:nowrap;
}
.rag-bar-wrap { width:100%; background:rgba(255,255,255,0.08); border-radius:4px; height:8px; overflow:hidden; }
.rag-bar-fill { height:100%; border-radius:4px; }
</style>
""", unsafe_allow_html=True)

    # ── Color scale strip ──────────────────────────────────────────────────────
    st.markdown('<div class="rag-section-label">RAG Color Scale</div>', unsafe_allow_html=True)
    segs = "".join(
        f'<div class="rag-strip-seg" style="background:#{s["hex"]};">'
        f'<span class="rag-strip-label" style="color:{_rag_text("#" + s["hex"])};">{s["label"]}</span>'
        f'<span class="rag-strip-hex"   style="color:{_rag_text("#" + s["hex"])};">#{s["hex"]}</span>'
        f'</div>'
        for s in RAG_SCALE
    )
    st.markdown(f'<div class="rag-strip">{segs}</div>', unsafe_allow_html=True)

    # ── Interactive gradient slider ────────────────────────────────────────────
    st.markdown('<div class="rag-section-label">Gradient Slider</div>', unsafe_allow_html=True)
    score = st.slider("RAG Score", 0, 100, 65, label_visibility="collapsed")
    pct   = score / 100
    color = _rag_color(float(score))
    tcolor = _rag_text(color)
    label = next((s["label"] for s in reversed(RAG_SCALE) if score >= s["score"]), RAG_SCALE[0]["label"])

    st.markdown(
        f'<div class="rag-grad-bar">'
        f'<div class="rag-grad-thumb" style="left:{pct*100:.1f}%;"></div>'
        f'</div>'
        f'<div class="rag-grad-labels"><span>Resistance · #EA3323</span><span>Momentum · #FFB547</span><span>Anchor · #00C27A</span></div>',
        unsafe_allow_html=True,
    )
    st.markdown(
        f'<span class="rag-badge" style="background:{color};color:{tcolor};">'
        f'{label} — {score}</span>',
        unsafe_allow_html=True,
    )

    # ── Sample data grid ───────────────────────────────────────────────────────
    st.markdown('<div class="rag-section-label" style="margin-top:2rem;">Sample Grid</div>', unsafe_allow_html=True)

    SAMPLE_ROWS = [
        ("Enterprise Sales",    "Q2 Target",  82, "$4.2M",  "↑ +12%"),
        ("Mid-Market",          "Q2 Target",  61, "$1.8M",  "↔ Flat"),
        ("SMB",                 "Q2 Target",  38, "$0.6M",  "↓ -8%"),
        ("APAC Expansion",      "Milestone",  91, "Phase 3","↑ Ahead"),
        ("Product Adoption",    "NPS Goal",   74, "68 NPS", "↑ +4pt"),
        ("Renewal Rate",        "Retention",  55, "81%",    "↓ -3pt"),
        ("Talent Pipeline",     "Headcount",  22, "4 / 18", "↓ Critical"),
        ("Compliance Training", "Completion", 97, "97%",    "↑ On track"),
    ]

    header = (
        "<table class='rag-table'>"
        "<thead><tr>"
        "<th>Initiative</th><th>Metric</th><th>Score</th>"
        "<th>Progress</th><th>Value</th><th>Trend</th>"
        "</tr></thead><tbody>"
    )
    rows_html = ""
    for init, metric, sc, val, trend in SAMPLE_ROWS:
        c  = _rag_color(sc)
        tc = _rag_text(c)
        rows_html += (
            f"<tr>"
            f"<td>{init}</td>"
            f"<td style='color:rgba(255,255,255,0.45);font-size:11px;'>{metric}</td>"
            f"<td><span class='rag-badge' style='background:{c};color:{tc};'>{sc}</span></td>"
            f"<td style='width:120px;'>"
            f"<div class='rag-bar-wrap'>"
            f"<div class='rag-bar-fill' style='width:{sc}%;background:{c};'></div>"
            f"</div></td>"
            f"<td style='font-family:var(--fq-font-mono);font-size:11px;'>{val}</td>"
            f"<td style='font-size:11px;color:rgba(255,255,255,0.55);'>{trend}</td>"
            f"</tr>"
        )
    st.markdown(header + rows_html + "</tbody></table>", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# TAB 7 — SEMANTIC INGEST
# ══════════════════════════════════════════════════════════════════════════════
with tab_semantic:
    if _HAS_SEMANTIC_INGEST:
        render_semantic_ingest_tab()
    else:
        st.error(
            "Semantic Ingest tab failed to load. "
            "Ensure `services/` and `tabs/` are present in the project directory."
        )


# ══════════════════════════════════════════════════════════════════════════════
# TAB 8 — GUIDED CONVERT
# Upload → rasterize → user designates cover/divider/ending per slide →
# convert() with forced_layouts → download output PPTX
# ══════════════════════════════════════════════════════════════════════════════
with tab_guided:
    st.subheader("Guided Convert")
    st.caption(
        "Upload a deck, review slide thumbnails, designate cover / divider / ending slides, "
        "then convert. All other slides default to Title Only or Sub layout."
    )

    _gc_upload = st.file_uploader(
        "Upload source PPTX", type=["pptx"], key="gc_upload"
    )

    if _gc_upload is not None:
        import hashlib as _hashlib

        _gc_hash = _hashlib.md5(_gc_upload.getvalue()).hexdigest()

        # ── Rasterize once per file (cached in session state) ─────────────────
        if st.session_state.get("gc_hash") != _gc_hash:
            with st.spinner("Rasterizing slides…"):
                try:
                    import tempfile as _gc_tf
                    from slide_vision import rasterize_pptx as _gc_rasterize

                    _gc_tmp_dir = _gc_tf.mkdtemp(prefix="gc_")
                    _gc_src = Path(_gc_tmp_dir) / _gc_upload.name
                    _gc_src.write_bytes(_gc_upload.getvalue())
                    _gc_pngs = _gc_rasterize(_gc_src, Path(_gc_tmp_dir))

                    st.session_state["gc_hash"]    = _gc_hash
                    st.session_state["gc_pngs"]    = [p.read_bytes() for p in _gc_pngs]
                    st.session_state["gc_src_bytes"] = _gc_upload.getvalue()
                    st.session_state["gc_src_name"]  = _gc_upload.name
                    st.session_state["gc_n"]       = len(_gc_pngs)
                except Exception as _gc_re:
                    st.error(f"Could not generate slide previews: {_gc_re}")
                    st.stop()

        _gc_n        = st.session_state.get("gc_n", 0)
        _gc_png_list = st.session_state.get("gc_pngs", [])

        if not _gc_png_list:
            st.warning("No slide previews available.")
        else:
            st.markdown(f"**{_gc_n} slides** — select a type for each:")

            # ── Thumbnail grid (3 columns) ────────────────────────────────────
            _GC_COLS = 3
            _gc_designations = {}

            for _gc_row_start in range(0, _gc_n, _GC_COLS):
                _gc_cols = st.columns(_GC_COLS)
                for _gc_col_idx in range(_GC_COLS):
                    _gc_slide_idx = _gc_row_start + _gc_col_idx + 1
                    if _gc_slide_idx > _gc_n:
                        break
                    with _gc_cols[_gc_col_idx]:
                        st.image(
                            _gc_png_list[_gc_slide_idx - 1],
                            caption=f"Slide {_gc_slide_idx}",
                            use_container_width=True,
                        )
                        # Default: slide 1 → cover, last slide → ending, rest → content
                        if _gc_slide_idx == 1:
                            _gc_default = "cover"
                        elif _gc_slide_idx == _gc_n:
                            _gc_default = "ending"
                        else:
                            _gc_default = "content"

                        _gc_designations[_gc_slide_idx] = st.selectbox(
                            f"slide_{_gc_slide_idx}",
                            ["content", "cover", "divider", "ending"],
                            index=["content", "cover", "divider", "ending"].index(_gc_default),
                            key=f"gc_des_{_gc_slide_idx}",
                            label_visibility="collapsed",
                        )

            # ── Validation ────────────────────────────────────────────────────
            _gc_covers  = [i for i, d in _gc_designations.items() if d == "cover"]
            _gc_endings = [i for i, d in _gc_designations.items() if d == "ending"]
            _gc_valid   = True

            if len(_gc_covers) > 1:
                st.warning(f"Only 1 cover allowed — slides {_gc_covers} are all marked as cover.")
                _gc_valid = False
            if len(_gc_endings) > 1:
                st.warning(f"Only 1 ending allowed — slides {_gc_endings} are all marked as ending.")
                _gc_valid = False

            # ── Convert button ────────────────────────────────────────────────
            if st.button("Convert", type="primary", key="gc_convert", disabled=not _gc_valid):
                # Build forced_layouts — only non-content slides get forced
                _gc_layout_map = {
                    "cover":   "Cover_Light",
                    "divider": "Divider_Dark",
                    "ending":  "Ending_PivotPurple",
                }
                _gc_forced = {
                    idx: _gc_layout_map[des]
                    for idx, des in _gc_designations.items()
                    if des in _gc_layout_map
                }

                with st.spinner("Converting…"):
                    try:
                        import tempfile as _gc_tf2
                        import convert_deck as _gc_cd

                        _gc_tmp2   = _gc_tf2.mkdtemp(prefix="gc_out_")
                        _gc_src2   = Path(_gc_tmp2) / st.session_state["gc_src_name"]
                        _gc_src2.write_bytes(st.session_state["gc_src_bytes"])

                        # Reload module to pick up any live edits
                        import importlib as _gc_il
                        _gc_il.reload(_gc_cd)

                        _gc_out = _gc_cd.convert(_gc_src2, forced_layouts=_gc_forced)
                        _gc_out_bytes = _gc_out.read_bytes()

                        st.success(f"Done — {len(_gc_out_bytes):,} bytes")
                        st.download_button(
                            label="⬇ Download converted PPTX",
                            data=_gc_out_bytes,
                            file_name=_gc_out.name,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            key="gc_download",
                        )

                        # Layout summary table
                        st.markdown("**Layout selection summary:**")
                        _gc_summary = []
                        for _si, _des in _gc_designations.items():
                            _forced_name = _gc_forced.get(_si, "—")
                            _gc_summary.append(
                                f"| {_si} | {_des} | {_forced_name} |"
                            )
                        st.markdown(
                            "| Slide | Designation | Forced Layout |\n"
                            "|---|---|---|\n" + "\n".join(_gc_summary)
                        )

                    except Exception as _gc_e:
                        st.error(f"Conversion failed: {_gc_e}")
                        import traceback as _gc_tb
                        st.code(_gc_tb.format_exc())
